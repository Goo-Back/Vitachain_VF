"""Présentation d'avancement VitaChain — version technique détaillée.

3 axes :
  01 — Base de données
  02 — Backend & Frontend
  03 — Sécurité

Chaque axe contient au moins 2 « preuves techniques » (extraits de code réels
tirés du repo : SQL PL/pgSQL, Python FastAPI, TypeScript Next.js).
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------- Palette ----------
NAVY      = RGBColor(0x0B, 0x1F, 0x3A)
TEAL      = RGBColor(0x00, 0xB8, 0xA9)
GOLD      = RGBColor(0xF5, 0xC1, 0x42)
LIGHT     = RGBColor(0xF4, 0xF7, 0xFA)
GREY      = RGBColor(0x55, 0x60, 0x6B)
GREY_LIGHT= RGBColor(0xDD, 0xE2, 0xE8)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
RED       = RGBColor(0xE0, 0x49, 0x4B)
GREEN     = RGBColor(0x2E, 0xA4, 0x6E)
CODE_BG   = RGBColor(0x12, 0x1A, 0x2B)
CODE_KW   = RGBColor(0xF5, 0xC1, 0x42)   # mots-clés
CODE_STR  = RGBColor(0x9B, 0xE3, 0x88)   # strings
CODE_COM  = RGBColor(0x7A, 0x87, 0x97)   # commentaires

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]


# ============================================================
# Helpers
# ============================================================
def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=NAVY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri",
             italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=12, color=NAVY, bullet="•"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(3)
        r = p.add_run()
        r.text = f"{bullet}  {item}"
        r.font.name = "Calibri"
        r.font.size = Pt(size)
        r.font.color.rgb = color


def code_box(slide, x, y, w, h, lines, *, lang_hint="sql", title=None,
             size=10):
    """Boîte de code avec coloration légère (mots-clés/commentaires)."""
    add_rect(slide, x, y, w, h, CODE_BG)
    if title:
        add_text(slide, x + Inches(0.12), y + Inches(0.05),
                 w - Inches(0.24), Inches(0.3),
                 title, size=10, bold=True, color=CODE_KW,
                 font="Consolas", italic=True)
        body_y = y + Inches(0.32)
        body_h = h - Inches(0.34)
    else:
        body_y = y + Inches(0.1)
        body_h = h - Inches(0.15)

    tb = slide.shapes.add_textbox(x + Inches(0.12), body_y,
                                  w - Inches(0.24), body_h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)

    KW = {
        "sql": {"create", "function", "language", "security", "definer", "as",
                "begin", "end", "if", "then", "else", "return", "returns",
                "declare", "select", "from", "where", "and", "or", "not",
                "exists", "raise", "exception", "using", "alter", "table",
                "enable", "row", "level", "policy", "for", "with", "check",
                "to", "grant", "revoke", "execute", "on", "in", "join",
                "trigger", "event", "default", "primary", "key", "references",
                "constraint", "check", "unique", "index", "where", "insert",
                "update", "delete", "values", "do", "perform"},
        "py":  {"def", "async", "await", "return", "if", "else", "elif",
                "try", "except", "raise", "from", "import", "as", "with",
                "for", "in", "not", "is", "None", "True", "False", "class",
                "lambda", "self", "and", "or", "pass"},
        "ts":  {"const", "let", "var", "function", "return", "if", "else",
                "async", "await", "import", "from", "export", "default",
                "class", "new", "this", "type", "interface", "extends",
                "implements", "true", "false", "null", "undefined", "try",
                "catch", "throw", "for", "of", "in"},
    }.get(lang_hint, set())

    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(0)
        # commentaire
        stripped = line.lstrip()
        if (stripped.startswith("--") or stripped.startswith("#")
                or stripped.startswith("//")):
            r = p.add_run(); r.text = line
            r.font.name = "Consolas"; r.font.size = Pt(size)
            r.font.color.rgb = CODE_COM
            continue
        # split simple par espaces tout en conservant l'indentation
        i = 0
        # indent
        indent = len(line) - len(line.lstrip())
        if indent:
            r = p.add_run(); r.text = " " * indent
            r.font.name = "Consolas"; r.font.size = Pt(size)
            r.font.color.rgb = WHITE
        remaining = line[indent:]
        # tokenisation très simple : espaces séparateurs
        buf = ""
        in_str = False
        str_char = ""
        for ch in remaining:
            if in_str:
                buf += ch
                if ch == str_char:
                    r = p.add_run(); r.text = buf
                    r.font.name = "Consolas"; r.font.size = Pt(size)
                    r.font.color.rgb = CODE_STR
                    buf = ""; in_str = False
                continue
            if ch in ("'", '"'):
                if buf:
                    r = p.add_run(); r.text = buf
                    r.font.name = "Consolas"; r.font.size = Pt(size)
                    r.font.color.rgb = WHITE
                    buf = ""
                in_str = True; str_char = ch; buf = ch
                continue
            if ch.isspace() or ch in "(),;:=<>+/*[]{}":
                if buf:
                    is_kw = buf.lower() in KW
                    r = p.add_run(); r.text = buf
                    r.font.name = "Consolas"; r.font.size = Pt(size)
                    r.font.color.rgb = CODE_KW if is_kw else WHITE
                    buf = ""
                r = p.add_run(); r.text = ch
                r.font.name = "Consolas"; r.font.size = Pt(size)
                r.font.color.rgb = WHITE
            else:
                buf += ch
        if buf:
            is_kw = buf.lower() in KW
            r = p.add_run(); r.text = buf
            r.font.name = "Consolas"; r.font.size = Pt(size)
            r.font.color.rgb = CODE_KW if is_kw else WHITE


def header_band(slide, title, subtitle=None, axis=None):
    add_rect(slide, 0, 0, SW, Inches(0.95), NAVY)
    add_rect(slide, 0, Inches(0.95), SW, Inches(0.06), TEAL)
    add_text(slide, Inches(0.5), Inches(0.15), Inches(11), Inches(0.5),
             title, size=26, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.58), Inches(11), Inches(0.4),
                 subtitle, size=12, color=GOLD, italic=True)
    add_text(slide, Inches(11.4), Inches(0.25), Inches(1.7), Inches(0.4),
             "VitaChain", size=15, bold=True, color=TEAL, align=PP_ALIGN.RIGHT)
    if axis:
        add_rect(slide, Inches(11.4), Inches(0.6), Inches(1.7), Inches(0.3),
                 GOLD)
        add_text(slide, Inches(11.4), Inches(0.62), Inches(1.7), Inches(0.3),
                 axis, size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


def footer(slide, page, section=""):
    add_rect(slide, 0, Inches(7.22), SW, Inches(0.28), NAVY)
    add_text(slide, Inches(0.4), Inches(7.24), Inches(8), Inches(0.26),
             f"VitaChain — État d'avancement | Mai 2026  •  {section}",
             size=9, color=WHITE)
    add_text(slide, Inches(11.5), Inches(7.24), Inches(1.4), Inches(0.26),
             f"{page} / 25", size=9, color=WHITE, align=PP_ALIGN.RIGHT)


def section_divider(num, title, subtitle, idx):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, NAVY)
    add_rect(s, 0, Inches(3.3), SW, Inches(0.06), TEAL)
    add_text(s, Inches(0.8), Inches(1.9), Inches(3), Inches(1.4),
             num, size=120, bold=True, color=TEAL)
    add_text(s, Inches(3.5), Inches(2.4), Inches(9), Inches(1),
             title, size=42, bold=True, color=WHITE)
    add_text(s, Inches(3.5), Inches(3.55), Inches(9), Inches(0.6),
             subtitle, size=17, color=GOLD, italic=True)
    # liste contenu de la section
    add_text(s, Inches(0.8), Inches(4.5), Inches(12), Inches(0.4),
             "Au programme de cette section :", size=14, bold=True, color=WHITE)
    return s


# ============================================================
# Slide 1 — Couverture
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, NAVY)
add_rect(s, 0, Inches(6.2), SW, Inches(0.15), TEAL)
add_rect(s, 0, Inches(6.4), SW, Inches(0.05), GOLD)

add_text(s, Inches(0.8), Inches(1.8), Inches(11.5), Inches(1.0),
         "VitaChain", size=64, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(2.75), Inches(11.5), Inches(0.55),
         "Plateforme agricole connectée — IoT, IA & Logistique",
         size=22, color=TEAL)
add_text(s, Inches(0.8), Inches(3.4), Inches(11.5), Inches(0.5),
         "État d'avancement technique — soutenance PFE",
         size=18, color=GOLD, italic=True)

# bandeau caractéristiques
add_rect(s, Inches(0.8), Inches(4.4), Inches(11.7), Inches(1.3), CODE_BG)
add_text(s, Inches(1.0), Inches(4.55), Inches(11.5), Inches(0.4),
         "Stack technique",
         size=12, bold=True, color=CODE_KW, italic=True, font="Consolas")
add_text(s, Inches(1.0), Inches(4.85), Inches(11.5), Inches(0.4),
         "PostgreSQL 15 + Supabase  •  FastAPI (Python 3.12)  •  "
         "Next.js 15 / React 19  •  Gemini 1.5",
         size=13, color=WHITE, font="Consolas")
add_text(s, Inches(1.0), Inches(5.2), Inches(11.5), Inches(0.4),
         "RLS partout  •  JWT custom claims  •  bcrypt devices  •  "
         "LISTEN/NOTIFY workers async",
         size=13, color=WHITE, font="Consolas")

add_text(s, Inches(0.8), Inches(6.7), Inches(11.5), Inches(0.4),
         "Yasser  —  PFE 2026", size=14, color=WHITE)


# ============================================================
# Slide 2 — Sommaire
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, LIGHT)
header_band(s, "Sommaire", "Plan de la présentation — ~25 minutes")

items = [
    ("00", "Contexte & architecture", "3 slides", TEAL,
     "Vision, stack, modules métier, métriques"),
    ("01", "Base de données", "5 slides", GOLD,
     "PostgreSQL + Supabase — RLS forcée, JWT hook, NOTIFY, anonymisation"),
    ("02", "Backend & Frontend", "7 slides", NAVY,
     "FastAPI + 7 workers • Next.js 15 + middleware d'auth • API contract"),
    ("03", "Sécurité", "5 slides", RED,
     "Défense en profondeur • bcrypt devices • surface d'attaque • garde-fous"),
    ("04", "Démo & roadmap", "3 slides", GREEN,
     "Démo live (KAT/FAR) • prochaines étapes • Q&A"),
]
y = Inches(1.4)
for num, title, count, color, desc in items:
    add_rect(s, Inches(0.5), y, Inches(0.9), Inches(1.0), color)
    add_text(s, Inches(0.5), y + Inches(0.2), Inches(0.9), Inches(0.6),
             num, size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.6), y + Inches(0.1), Inches(8), Inches(0.4),
             title, size=18, bold=True, color=NAVY)
    add_text(s, Inches(1.6), y + Inches(0.55), Inches(10), Inches(0.4),
             desc, size=11, color=GREY)
    add_rect(s, Inches(11.5), y + Inches(0.3), Inches(1.4), Inches(0.4),
             LIGHT, line=color)
    add_text(s, Inches(11.5), y + Inches(0.33), Inches(1.4), Inches(0.4),
             count, size=11, bold=True, color=color, align=PP_ALIGN.CENTER)
    y += Inches(1.13)
footer(s, 2, "Sommaire")


# ============================================================
# Slide 3 — Architecture globale
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, LIGHT)
header_band(s, "Architecture globale", "Une stack 100% Supabase + FastAPI")

layers = [
    ("Devices  /  Clients",
     "ESP32 firmware C++  •  Navigateur (desktop + mobile)  •  Webhooks externes",
     GREY),
    ("Frontend  —  Next.js 15 / React 19",
     "App Router (SSR + RSC)  •  Tailwind 4  •  @supabase/ssr  •  Leaflet  •  Sentry",
     TEAL),
    ("Backend  —  FastAPI + Workers asynchrones",
     "16 routers REST  •  7 workers (asyncpg LISTEN)  •  Gemini  •  Brevo  •  OWM  •  Sentinel-2",
     GOLD),
    ("Base de données  —  PostgreSQL 15 (Supabase)",
     "40 migrations  •  54 policies RLS  •  Auth Hook PL/pgSQL  •  Storage  •  Realtime",
     NAVY),
]
y = Inches(1.4)
for title, sub, color in layers:
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(1.2), WHITE,
             line=GREY_LIGHT)
    add_rect(s, Inches(0.5), y, Inches(0.25), Inches(1.2), color)
    add_text(s, Inches(0.95), y + Inches(0.18), Inches(11.5), Inches(0.45),
             title, size=17, bold=True, color=NAVY)
    add_text(s, Inches(0.95), y + Inches(0.65), Inches(11.5), Inches(0.4),
             sub, size=11, color=GREY, font="Consolas")
    y += Inches(1.35)
footer(s, 3, "Architecture")


# ============================================================
# Slide 4 — Modules métier + métriques
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, LIGHT)
header_band(s, "Modules métier & métriques",
            "Architecture modulaire — un domaine = un module backend")

# 4 modules à gauche
mods = [
    ("AUTH", "Authentification & KYC",   TEAL),
    ("KAT",  "Katara IoT (champ)",       GOLD),
    ("FAR",  "FarMarket (marketplace)",  RED),
    ("SCN",  "SecondServe (à venir)",    GREY),
]
y = Inches(1.4)
for code, name, color in mods:
    add_rect(s, Inches(0.5), y, Inches(1.0), Inches(0.95), color)
    add_text(s, Inches(0.5), y + Inches(0.2), Inches(1.0), Inches(0.55),
             code, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.65), y + Inches(0.3), Inches(5.5), Inches(0.4),
             name, size=15, bold=True, color=NAVY)
    y += Inches(1.05)

# KPI cards à droite
kpis = [
    ("40",  "Migrations SQL\nversionnées",    TEAL),
    ("54",  "Policies RLS\n(100% des tables)",GOLD),
    ("63",  "Fonctions /\ntriggers SQL",      NAVY),
    ("74",  "Fichiers Python\n(backend)",     RED),
    ("25",  "Pages\n(Next.js)",               GREEN),
    ("202", "Fichiers de tests\n(back + front)", TEAL),
]
positions = [
    (Inches(7.4), Inches(1.4)),
    (Inches(9.3), Inches(1.4)),
    (Inches(11.2), Inches(1.4)),
    (Inches(7.4), Inches(3.4)),
    (Inches(9.3), Inches(3.4)),
    (Inches(11.2), Inches(3.4)),
]
for (val, lbl, color), (x, y) in zip(kpis, positions):
    add_rect(s, x, y, Inches(1.8), Inches(1.85), color)
    add_text(s, x, y + Inches(0.15), Inches(1.8), Inches(0.8),
             val, size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(1.05), Inches(1.8), Inches(0.7),
             lbl, size=10, color=WHITE, align=PP_ALIGN.CENTER)

# encart bas
add_rect(s, Inches(0.5), Inches(5.55), Inches(12.3), Inches(1.4),
         WHITE, line=GREY_LIGHT)
add_rect(s, Inches(0.5), Inches(5.55), Inches(0.15), Inches(1.4), GOLD)
add_text(s, Inches(0.8), Inches(5.65), Inches(11.7), Inches(0.4),
         "Principe directeur : un module = un domaine, isolé via des sous-paths "
         "et des préfixes de tables (m1_, m2_).", size=12, bold=True, color=NAVY)
add_text(s, Inches(0.8), Inches(6.05), Inches(11.7), Inches(0.9),
         "→ Chaque module embarque : ses migrations SQL, son router FastAPI, "
         "ses workers asynchrones, ses pages frontend, ses tests.\n"
         "→ Couplage faible entre modules — un module peut être désactivé "
         "(SCN) sans casser les autres (KAT, FAR).",
         size=11, color=GREY)
footer(s, 4, "Architecture")


# ============================================================
# SECTION 01 — BASE DE DONNÉES
# ============================================================
s = section_divider("01", "Base de données",
                    "PostgreSQL 15 / Supabase — schéma versionné, RLS partout", 1)
items = [
    "Migrations versionnées (40 fichiers, idempotents, rejouables)",
    "Modèle de données — extrait des 25+ tables métier",
    "PREUVE 1 — Event trigger qui REFUSE toute table sans RLS",
    "PREUVE 2 — JWT Auth Hook PL/pgSQL (rôle injecté à la connexion)",
    "Flux temps réel (LISTEN/NOTIFY) + anonymisation FarMarket",
]
y = Inches(4.95)
for it in items:
    add_text(s, Inches(0.8), y, Inches(12), Inches(0.4),
             "›  " + it, size=14, color=WHITE)
    y += Inches(0.35)


# ============================================================
# Slide 6 — Migrations versionnées
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, LIGHT)
header_band(s, "Migrations versionnées", "Un fichier .sql = un changement atomique",
            axis="AXE 01")

# Timeline horizontale
add_text(s, Inches(0.5), Inches(1.3), Inches(12), Inches(0.4),
         "Chronologie — 40 migrations groupées par épopée",
         size=14, bold=True, color=NAVY)

epics_tl = [
    ("INF / AUTH", "0001 → 0015", "15 mig.", TEAL,
     "Profiles, JWT hook, KYC, RLS contract, Storage policies"),
    ("KAT (Katara)", "0016 → 0031", "16 mig.", GOLD,
     "Parcelles, devices, télémétrie, seuils, diagnostics IA, history"),
    ("FAR (FarMarket)", "0032 → 0042", "11 mig.", RED,
     "Ads, photos, pivot logistique, orders, tracking, notifications"),
]
y = Inches(1.85)
for name, range_, count, color, desc in epics_tl:
    add_rect(s, Inches(0.5), y, Inches(2.5), Inches(0.9), color)
    add_text(s, Inches(0.5), y + Inches(0.13), Inches(2.5), Inches(0.4),
             name, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), y + Inches(0.5), Inches(2.5), Inches(0.4),
             count, size=11, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(3.15), y + Inches(0.05), Inches(2), Inches(0.4),
             range_, size=14, bold=True, color=NAVY, font="Consolas")
    add_text(s, Inches(3.15), y + Inches(0.45), Inches(9.5), Inches(0.4),
             desc, size=11, color=GREY)
    y += Inches(1.05)

# Bonnes pratiques
add_rect(s, Inches(0.5), Inches(5.0), Inches(6.0), Inches(1.95),
         WHITE, line=GREY_LIGHT)
add_rect(s, Inches(0.5), Inches(5.0), Inches(6.0), Inches(0.4), NAVY)
add_text(s, Inches(0.5), Inches(5.05), Inches(6.0), Inches(0.3),
         "  Pratiques appliquées", size=12, bold=True, color=WHITE)
add_bullets(s, Inches(0.65), Inches(5.5), Inches(5.8), Inches(1.4), [
    "Idempotence : if not exists, create or replace partout",
    "Bookkeeping : table _migrations rejoue sans casser",
    "Format de nom : NNNN_<epic>_<short_desc>.sql",
    "Commentaire d'en-tête : story, raison, contraintes",
], size=11, color=GREY)

add_rect(s, Inches(6.8), Inches(5.0), Inches(6.0), Inches(1.95),
         WHITE, line=GREY_LIGHT)
add_rect(s, Inches(6.8), Inches(5.0), Inches(6.0), Inches(0.4), TEAL)
add_text(s, Inches(6.8), Inches(5.05), Inches(6.0), Inches(0.3),
         "  Garde-fous CI", size=12, bold=True, color=WHITE)
add_bullets(s, Inches(6.95), Inches(5.5), Inches(5.8), Inches(1.4), [
    "scripts/verify-rls-enabled.sh : refuse table sans RLS",
    "db/tests/auth04_rls_contract.sql : assertions pgSQL",
    "Event trigger trg_enforce_rls_on_public_tables (DB)",
    "Test allowlist : interdit l'usage non documenté de service_role",
], size=11, color=GREY)
footer(s, 6, "01 — Base de données")


# ============================================================
# Slide 7 — Modèle de données
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, LIGHT)
header_band(s, "Modèle de données — extrait",
            "Relations principales et choix techniques", axis="AXE 01")

# Gauche : tables
left_x = Inches(0.5)
add_rect(s, left_x, Inches(1.3), Inches(6.0), Inches(5.7), WHITE,
         line=GREY_LIGHT)
add_rect(s, left_x, Inches(1.3), Inches(6.0), Inches(0.45), NAVY)
add_text(s, left_x, Inches(1.36), Inches(6.0), Inches(0.35),
         "Tables principales", size=13, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)

tbls = [
    ("profiles",                    "AUTH", "id, role, verification_status"),
    ("kyc_documents",               "AUTH", "doc_type, status, reviewer"),
    ("m1_katara_parcels",           "KAT",  "geojson polygon, farmer_id"),
    ("m1_katara_devices",           "KAT",  "api_key_hash (bcrypt), parcel_id"),
    ("m1_katara_telemetry",         "KAT",  "soil_ph, EC, moisture, recorded_at"),
    ("m1_katara_telemetry_history", "KAT",  "agrégat horaire (provenance KAT-13)"),
    ("m1_katara_ai_diagnostics",    "KAT",  "prompt, gemini_response, NDVI, status"),
    ("m1_katara_alert_thresholds",  "KAT",  "min/max par mesure, par parcelle"),
    ("m2_farmarket_ads",            "FAR",  "produit, prix MAD, photos[]"),
    ("m2_farmarket_orders",         "FAR",  "8 statuts, restaurant_id, totals"),
    ("m2_farmarket_order_items",    "FAR",  "snapshot prix, farmer_id, statut"),
    ("notifications_outbox",        "CORE", "pattern outbox → Brevo email"),
]
y = Inches(1.85)
for name, dom, desc in tbls:
    add_text(s, left_x + Inches(0.15), y, Inches(2.7), Inches(0.4),
             name, size=10, bold=True, color=NAVY, font="Consolas")
    add_text(s, left_x + Inches(2.8), y, Inches(0.6), Inches(0.4),
             dom, size=10, bold=True, color=TEAL)
    add_text(s, left_x + Inches(3.4), y, Inches(2.6), Inches(0.4),
             desc, size=10, color=GREY)
    y += Inches(0.42)

# Droite : choix techniques
right_x = Inches(6.8)
add_rect(s, right_x, Inches(1.3), Inches(6.0), Inches(5.7), WHITE,
         line=GREY_LIGHT)
add_rect(s, right_x, Inches(1.3), Inches(6.0), Inches(0.45), TEAL)
add_text(s, right_x, Inches(1.36), Inches(6.0), Inches(0.35),
         "Choix techniques notables", size=13, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)

choices = [
    "Préfixes m1_ / m2_ pour scoper par module (Katara, FarMarket)",
    "Enums Postgres natifs (device_status, order_status) — pas de varchar libre",
    "Index uniques PARTIELS — ex: 1 device actif par parcelle (where status<>'UNLINKED')",
    "Hypertable-like sur telemetry — agrégation horaire dans table history dédiée",
    "Triggers BEFORE INSERT pour dénormalisation contrôlée (farmer_id sync)",
    "NOTIFY canaux : telemetry_inserted, threshold_breach, order_status_changed",
    "Vues SECURITY INVOKER → respect de la RLS de l'appelant",
    "Snapshot pricing sur order_items : immune aux edits d'ads après commande",
    "Audit guards SQL : RAISE EXCEPTION si INSERT direct (force passage par RPC)",
    "Format clé device : vk_<32 hex> — 128 bits d'entropie, hashé bcrypt cost=10",
]
add_bullets(s, right_x + Inches(0.2), Inches(1.85), Inches(5.6), Inches(5.0),
            choices, size=11, color=GREY)
footer(s, 7, "01 — Base de données")


# ============================================================
# Slide 8 — PREUVE 1 — Event trigger RLS
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, LIGHT)
header_band(s, "Preuve 1 — Event trigger : RLS obligatoire",
            "migration 0009 — il est impossible de créer une table sans RLS",
            axis="AXE 01")

# Explication gauche
add_rect(s, Inches(0.5), Inches(1.3), Inches(5.0), Inches(5.7),
         WHITE, line=GREY_LIGHT)
add_rect(s, Inches(0.5), Inches(1.3), Inches(0.12), Inches(5.7), RED)
add_text(s, Inches(0.75), Inches(1.45), Inches(4.7), Inches(0.4),
         "Le problème", size=14, bold=True, color=NAVY)
add_bullets(s, Inches(0.75), Inches(1.85), Inches(4.7), Inches(1.5), [
    "Une table créée sans ENABLE RLS = données exposées",
    "Le dev peut oublier → leak silencieux",
    "Les outils CI/Shell scripts peuvent passer à côté",
], size=11, color=GREY)

add_text(s, Inches(0.75), Inches(3.4), Inches(4.7), Inches(0.4),
         "La solution : 3 lignes de défense", size=14, bold=True, color=NAVY)
add_bullets(s, Inches(0.75), Inches(3.8), Inches(4.7), Inches(3), [
    "① Test pgSQL — auth04_rls_contract.sql",
    "② Shell CI — verify-rls-enabled.sh",
    "③ Event trigger Postgres ← DERNIÈRE LIGNE",
    "",
    "L'event trigger s'exécute DANS la base, quel que soit",
    "l'auteur du CREATE TABLE (migration, psql, dashboard).",
    "Si la RLS n'est pas activée → abort de la transaction.",
], size=11, color=GREY)

# Code SQL droite
code_lines = [
    "create event trigger trg_enforce_rls_on_public_tables",
    "    on ddl_command_end",
    "    when tag in ('CREATE TABLE')",
    "    execute function public.enforce_rls_on_public_tables();",
    "",
    "create or replace function public.enforce_rls_on_public_tables()",
    "returns event_trigger language plpgsql security definer as $$",
    "declare obj record; begin",
    "  for obj in select * from pg_event_trigger_ddl_commands()",
    "             where command_tag = 'CREATE TABLE' loop",
    "    if obj.object_identity like 'public.%' then",
    "      if not exists (",
    "        select 1 from pg_class c join pg_namespace n",
    "          on n.oid = c.relnamespace",
    "         where n.nspname = 'public'",
    "           and c.relname = split_part(obj.object_identity,'.',2)",
    "           and c.relrowsecurity = true",
    "      ) then",
    "        raise exception",
    "          'AUTH-04: table %s was created without RLS',",
    "          obj.object_identity using errcode = '42501';",
    "      end if;",
    "    end if;",
    "  end loop;",
    "end; $$;",
]
code_box(s, Inches(5.7), Inches(1.3), Inches(7.1), Inches(5.7),
         code_lines, lang_hint="sql",
         title="db/migrations/0009_auth04_force_rls_contract.sql", size=10)
footer(s, 8, "01 — Base de données")


# ============================================================
# Slide 9 — PREUVE 2 — JWT custom hook
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, LIGHT)
header_band(s, "Preuve 2 — JWT enrichi à la connexion",
            "migration 0006 — Auth Hook PL/pgSQL injecte le rôle dans le token",
            axis="AXE 01")

# Bandeau pourquoi
add_rect(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.9),
         CODE_BG)
add_text(s, Inches(0.7), Inches(1.4), Inches(12), Inches(0.4),
         "Pourquoi ?", size=13, bold=True, color=GOLD, font="Consolas")
add_text(s, Inches(0.7), Inches(1.75), Inches(12), Inches(0.45),
         "Sans claim → chaque policy RLS doit faire SELECT role FROM profiles "
         "(1 lookup/évaluation).  Avec claim → coût ZÉRO côté policy.",
         size=11, color=WHITE, font="Consolas")

# Code SQL
code_lines = [
    "create or replace function public.custom_access_token_hook(event jsonb)",
    "returns jsonb language plpgsql security definer stable",
    "set search_path = public, pg_temp",
    "as $$",
    "declare",
    "    uid           uuid;",
    "    resolved_role text;",
    "    new_claims    jsonb;",
    "begin",
    "    uid := (event->>'user_id')::uuid;",
    "",
    "    select role::text into resolved_role",
    "      from public.profiles where id = uid;",
    "",
    "    if resolved_role is null then",
    "        return event;             -- défensif : pas de profil → token brut",
    "    end if;",
    "",
    "    new_claims := coalesce(event->'claims', '{}'::jsonb)",
    "               || jsonb_build_object('user_role', resolved_role);",
    "",
    "    return jsonb_set(event, '{claims}', new_claims);",
    "end; $$;",
    "",
    "-- Verrouillage : seul l'admin Auth Supabase peut l'exécuter",
    "grant execute on function public.custom_access_token_hook(jsonb)",
    "    to supabase_auth_admin;",
    "revoke execute on function public.custom_access_token_hook(jsonb) from public;",
]
code_box(s, Inches(0.5), Inches(2.35), Inches(8.5), Inches(4.65),
         code_lines, lang_hint="sql",
         title="db/migrations/0006_auth02_jwt_role_hook.sql", size=10)

# Effet observable
add_rect(s, Inches(9.2), Inches(2.35), Inches(3.6), Inches(4.65),
         WHITE, line=GREY_LIGHT)
add_rect(s, Inches(9.2), Inches(2.35), Inches(3.6), Inches(0.4), GREEN)
add_text(s, Inches(9.2), Inches(2.4), Inches(3.6), Inches(0.3),
         "Effet observable", size=11, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(9.35), Inches(2.85), Inches(3.4), Inches(0.4),
         "Payload JWT après login :", size=10, bold=True, color=NAVY)
mini_code = [
    "{",
    '  "sub": "uuid-...",',
    '  "email": "f@x.fr",',
    '  "user_role": "FARMER",',
    '  "verification_status":',
    '       "VERIFIED",',
    '  "exp": 1748...',
    "}",
]
code_box(s, Inches(9.35), Inches(3.3), Inches(3.3), Inches(2.0),
         mini_code, lang_hint="ts", size=10)
add_text(s, Inches(9.35), Inches(5.45), Inches(3.4), Inches(1.5),
         "→ Policy RLS :\n   USING (has_role('FARMER')\n     AND auth.uid() = farmer_id)\n\n"
         "Zéro round-trip SQL\nsupplémentaire.",
         size=10, color=GREY, font="Consolas")
footer(s, 9, "01 — Base de données")


# ============================================================
# Slide 10 — Flux temps réel + anonymisation
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, LIGHT)
header_band(s, "Flux temps réel & anonymisation",
            "LISTEN/NOTIFY + projection view pour la RGPD", axis="AXE 01")

# Pipeline horizontal (haut)
add_text(s, Inches(0.5), Inches(1.3), Inches(12), Inches(0.4),
         "① Pipeline télémétrie temps réel (Katara)",
         size=13, bold=True, color=NAVY)
steps = [
    ("ESP32",       "POST /ingest", TEAL),
    ("FastAPI",     "RPC SECURITY\nDEFINER", GOLD),
    ("Postgres",    "INSERT +\nNOTIFY", NAVY),
    ("Worker",      "LISTEN\nasyncpg", RED),
    ("Notif",       "outbox →\nBrevo", GREEN),
]
x = Inches(0.5)
for label, desc, color in steps:
    add_rect(s, x, Inches(1.75), Inches(2.3), Inches(1.3), color)
    add_text(s, x, Inches(1.85), Inches(2.3), Inches(0.4),
             label, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(2.3), Inches(2.3), Inches(0.9),
             desc, size=10, color=WHITE, align=PP_ALIGN.CENTER, font="Consolas")
    if x < Inches(10.5):
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 x + Inches(2.3), Inches(2.18),
                                 Inches(0.2), Inches(0.45))
        arr.fill.solid(); arr.fill.fore_color.rgb = GREY
        arr.line.fill.background()
    x += Inches(2.5)

add_text(s, Inches(0.5), Inches(3.2), Inches(12.3), Inches(0.5),
         "Latence mesurée : < 100 ms entre INSERT et déclenchement worker "
         "(vs 1-5s en polling).",
         size=11, color=GREY, italic=True)

# Anonymisation
add_text(s, Inches(0.5), Inches(3.85), Inches(12), Inches(0.4),
         "② Anonymisation FarMarket (BR-F5)",
         size=13, bold=True, color=NAVY)
add_rect(s, Inches(0.5), Inches(4.25), Inches(6.0), Inches(2.7),
         WHITE, line=GREY_LIGHT)
add_rect(s, Inches(0.5), Inches(4.25), Inches(0.12), Inches(2.7), RED)
add_text(s, Inches(0.75), Inches(4.35), Inches(5.6), Inches(0.4),
         "Contrainte métier", size=12, bold=True, color=NAVY)
add_bullets(s, Inches(0.75), Inches(4.7), Inches(5.6), Inches(2.1), [
    "Le farmer ne doit PAS voir l'identité du restaurant",
    "Mais doit pouvoir traiter ses items de commande",
    "Solution : table orders SANS policy SELECT pour farmers",
    "+ vue v_farmer_incoming_items qui projette un opaque resto_handle",
], size=11, color=GREY)

# Code SQL
code_lines = [
    "create view v_farmer_incoming_items",
    "with (security_invoker = true) as",
    "select",
    "    oi.id, oi.order_id, oi.product_name,",
    "    oi.quantity, oi.line_total_mad,",
    "    encode(",
    "      digest(",
    "        o.restaurant_id::text || ':'",
    "          || oi.farmer_id::text,",
    "        'sha256'",
    "      ),",
    "    'hex') as resto_handle",
    "from m2_farmarket_order_items oi",
    "join m2_farmarket_orders o on o.id = oi.order_id;",
]
code_box(s, Inches(6.8), Inches(4.25), Inches(6.0), Inches(2.7),
         code_lines, lang_hint="sql",
         title="vue avec hash SHA-256 du couple (resto, farmer)", size=10)
footer(s, 10, "01 — Base de données")


# ============================================================
# SECTION 02 — BACKEND & FRONTEND
# ============================================================
s = section_divider("02", "Backend & Frontend",
                    "FastAPI + 7 workers asynchrones  |  Next.js 15 / React 19",
                    2)
items = [
    "Backend FastAPI — 16 routers, 7 workers asynchrones",
    "PREUVE 3 — Endpoint d'ingestion télémétrie < 50 ms (hot path)",
    "PREUVE 4 — Pipeline diagnostic IA Gemini (orchestrator séquentiel)",
    "Frontend Next.js 15 — App Router, SSR, RSC",
    "PREUVE 5 — Middleware d'auth + gate de vérification KYC",
    "Communication Front ↔ Back — 2 canaux complémentaires",
]
y = Inches(4.95)
for it in items:
    add_text(s, Inches(0.8), y, Inches(12), Inches(0.4),
             "›  " + it, size=14, color=WHITE)
    y += Inches(0.32)


# ============================================================
# Slide 12 — Backend stack + routers
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, LIGHT)
header_band(s, "Backend — FastAPI",
            "16 routers REST + 7 workers asynchrones", axis="AXE 02")

# Stack à gauche
add_rect(s, Inches(0.5), Inches(1.3), Inches(4.0), Inches(5.7),
         WHITE, line=GREY_LIGHT)
add_rect(s, Inches(0.5), Inches(1.3), Inches(4.0), Inches(0.4), NAVY)
add_text(s, Inches(0.5), Inches(1.35), Inches(4.0), Inches(0.3),
         "  Stack technique", size=12, bold=True, color=WHITE)
stack = [
    ("FastAPI",       "0.115"),
    ("Pydantic v2",   "validation stricte"),
    ("supabase-py",   "client PostgREST"),
    ("asyncpg",       "LISTEN/NOTIFY direct"),
    ("PyJWT",         "vérif JWT HS256"),
    ("bcrypt",        "hash clés device"),
    ("Jinja2",        "autoescape on"),
    ("google-genai",  "Gemini 1.5 Flash"),
    ("tifffile/numpy","NDVI Sentinel-2"),
    ("structlog",     "logs JSON"),
    ("sentry-sdk",    "erreurs prod"),
    ("orjson",        "JSON rapide"),
    ("gunicorn",      "+ uvicorn workers"),
]
y = Inches(1.8)
for n, d in stack:
    add_text(s, Inches(0.65), y, Inches(1.7), Inches(0.3),
             n, size=10, bold=True, color=TEAL, font="Consolas")
    add_text(s, Inches(2.4), y, Inches(2.0), Inches(0.3),
             d, size=10, color=GREY)
    y += Inches(0.37)

# Routers au centre
add_rect(s, Inches(4.7), Inches(1.3), Inches(4.0), Inches(5.7),
         WHITE, line=GREY_LIGHT)
add_rect(s, Inches(4.7), Inches(1.3), Inches(4.0), Inches(0.4), GOLD)
add_text(s, Inches(4.7), Inches(1.35), Inches(4.0), Inches(0.3),
         "  Routers FastAPI (16)", size=12, bold=True, color=WHITE)
routers = [
    "/api/v1/health",
    "/api/v1/kyc",
    "/api/v1/admin/kyc",
    "/api/v1/admin/farmarket",
    "/api/v1/katara/parcels",
    "/api/v1/katara/parcels/.../devices",
    "/api/v1/katara/devices/.../unlink",
    "/api/v1/katara/ingest",
    "/api/v1/katara/telemetry",
    "/api/v1/katara/devices/.../history",
    "/api/v1/katara/thresholds",
    "/api/v1/katara/diagnostics",
    "/api/v1/katara/overview",
    "/api/v1/farmarket",
    "/api/v1/secondserve",
    "/api/v1/notifications",
]
y = Inches(1.8)
for r in routers:
    add_text(s, Inches(4.85), y, Inches(3.8), Inches(0.3),
             "› " + r, size=9, color=NAVY, font="Consolas")
    y += Inches(0.31)

# Workers à droite
add_rect(s, Inches(8.9), Inches(1.3), Inches(3.9), Inches(5.7),
         WHITE, line=GREY_LIGHT)
add_rect(s, Inches(8.9), Inches(1.3), Inches(3.9), Inches(0.4), RED)
add_text(s, Inches(8.9), Inches(1.35), Inches(3.9), Inches(0.3),
         "  Workers asynchrones (7)", size=12, bold=True, color=WHITE)

workers = [
    ("katara_threshold",         "LISTEN télémétrie → alertes"),
    ("katara_offline",           "détection silence > 1h"),
    ("katara_diagnostic",        "Gemini + OWM + NDVI"),
    ("katara_diagnostic_email",  "Markdown → HTML → email"),
    ("farmarket_expiry",         "expire ads périmées"),
    ("farmarket_order_notify",   "transitions commande"),
    ("notifications_mailer",     "outbox → Brevo"),
]
y = Inches(1.85)
for n, d in workers:
    add_text(s, Inches(9.05), y, Inches(3.6), Inches(0.35),
             n, size=10, bold=True, color=NAVY, font="Consolas")
    add_text(s, Inches(9.05), y + Inches(0.32), Inches(3.6), Inches(0.3),
             d, size=9, color=GREY)
    y += Inches(0.7)

add_text(s, Inches(9.05), Inches(6.7), Inches(3.6), Inches(0.3),
         "→ Process séparés, scalables", size=9, color=GREY, italic=True)

footer(s, 12, "02 — Backend & Frontend")


# ============================================================
# Slide 13 — PREUVE 3 — Ingest hot path
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, LIGHT)
header_band(s, "Preuve 3 — Endpoint d'ingestion télémétrie",
            "SLA : < 50 ms p50, < 150 ms p99 — un seul round-trip DB", axis="AXE 02")

# Code Python
code_lines = [
    "@router.post('', status_code=204, response_class=Response)",
    "async def ingest_telemetry(",
    "    payload: TelemetryPayload,",
    "    db: Annotated[Client, Depends(_ingest_db)],",
    "    x_device_id: str | None = Header(default=None),",
    "    x_device_api_key: str | None = Header(default=None),",
    ") -> Response:",
    "    if not x_device_id or not x_device_api_key:",
    "        raise HTTPException(401, detail=_INVALID_CREDS)",
    "",
    "    # Un seul appel DB :",
    "    # verify_bcrypt + insert + touch device + NOTIFY",
    "    res = db.rpc('m1_katara_ingest', {",
    "        'p_device_id_str':     x_device_id,",
    "        'p_api_key':           x_device_api_key,",
    "        'p_soil_moisture':     payload.soil_moisture,",
    "        'p_soil_temperature':  payload.soil_temperature,",
    "        'p_soil_ph':           payload.soil_ph,",
    "        'p_soil_conductivity': payload.soil_conductivity,",
    "        'p_battery_level':     payload.battery_level,",
    "        'p_recorded_at':       payload.recorded_at.isoformat(),",
    "    }).execute()",
    "",
    "    if res.data is None:",
    "        # Constant-time error : ne dit pas LEQUEL des deux",
    "        # (device_id ou api_key) est faux → anti-énumération",
    "        raise HTTPException(401, detail=_INVALID_CREDS)",
    "",
    "    return Response(status_code=204,",
    "        headers={'X-Telemetry-Id': str(res.data)})",
]
code_box(s, Inches(0.5), Inches(1.3), Inches(8.0), Inches(5.7),
         code_lines, lang_hint="py",
         title="backend/app/modules/katara/ingest.py", size=10)

# Explications à droite
add_rect(s, Inches(8.7), Inches(1.3), Inches(4.1), Inches(5.7),
         WHITE, line=GREY_LIGHT)
add_rect(s, Inches(8.7), Inches(1.3), Inches(0.12), Inches(5.7), TEAL)
add_text(s, Inches(8.95), Inches(1.4), Inches(3.7), Inches(0.4),
         "Décisions techniques", size=13, bold=True, color=NAVY)
add_bullets(s, Inches(8.95), Inches(1.85), Inches(3.7), Inches(2.5), [
    "Auth par device, pas par user (pas de JWT)",
    "Service-role client : seul autorisé à INSERT (RLS FORCE)",
    "RPC SQL bundlant tout → 1 round-trip",
    "204 No Content : ESP32 ne lit pas le body",
    "Headers X-Device-Id / X-Device-Api-Key",
], size=11, color=GREY)

add_text(s, Inches(8.95), Inches(4.5), Inches(3.7), Inches(0.4),
         "Sécurité défensive", size=13, bold=True, color=NAVY)
add_bullets(s, Inches(8.95), Inches(4.9), Inches(3.7), Inches(2.0), [
    "_INVALID_CREDS constant → pas d'énumération",
    "Bcrypt verify côté SQL (constant-time pgcrypto)",
    "Callsite ajouté à l'allowlist service_role",
    "Rate limit NGINX en amont (AUTH-08)",
], size=11, color=GREY)

footer(s, 13, "02 — Backend & Frontend")


# ============================================================
# Slide 14 — PREUVE 4 — Pipeline IA Gemini
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, LIGHT)
header_band(s, "Preuve 4 — Pipeline diagnostic IA Gemini",
            "Orchestrator séquentiel — taxonomie d'erreurs explicite",
            axis="AXE 02")

# Pipeline visuel haut
add_text(s, Inches(0.5), Inches(1.3), Inches(12), Inches(0.4),
         "Pipeline KAT-08 — 6 étapes, FAILED précis à chaque échec",
         size=13, bold=True, color=NAVY)
pipe = [
    ("Claim",     "FOR UPDATE\nSKIP LOCKED", TEAL),
    ("OWM",       "Météo 24h\nlat/lng", GOLD),
    ("NDVI",      "Sentinel-2\nTIFF float32", GREEN),
    ("Aggregate", "Moy. 7j\ntélémétrie", NAVY),
    ("Prompt",    "Jinja2\nautoescape", RED),
    ("Gemini",    "1.5 Flash\nstreaming", TEAL),
]
x = Inches(0.5)
for label, desc, color in pipe:
    add_rect(s, x, Inches(1.75), Inches(1.95), Inches(1.3), color)
    add_text(s, x, Inches(1.85), Inches(1.95), Inches(0.4),
             label, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(2.3), Inches(1.95), Inches(0.9),
             desc, size=9, color=WHITE, align=PP_ALIGN.CENTER, font="Consolas")
    if x < Inches(11.0):
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 x + Inches(1.95), Inches(2.15),
                                 Inches(0.15), Inches(0.5))
        arr.fill.solid(); arr.fill.fore_color.rgb = GREY
        arr.line.fill.background()
    x += Inches(2.1)

# Code Python — orchestrator
code_lines = [
    "# Chaque étape est try/except indépendant → FAILED précis",
    "try:",
    "    owm = await asyncio.to_thread(fetch_weather, lat, lng)",
    "except Exception as exc:",
    "    mark_failed(diag_id, f'owm_unavailable: {exc!r}')",
    "    return",
    "",
    "try:",
    "    ndvi = await asyncio.to_thread(fetch_ndvi, parcel_id, geojson)",
    "except Exception as exc:",
    "    mark_failed(diag_id, f'ndvi_unavailable: {exc!r}')",
    "    return",
    "",
    "# ... agrégat 7j, build prompt, appel Gemini ...",
    "",
    "try:",
    "    result_text = await call_gemini(prompt)",
    "except GeminiRateLimited as exc:",
    "    mark_failed(diag_id, f'gemini_rate_limited: {exc!r}')",
    "except GeminiUnavailable as exc:",
    "    mark_failed(diag_id, f'gemini_unavailable: {exc!r}')",
    "",
    "mark_completed(diag_id, result_text)   # NOTIFY → email worker",
]
code_box(s, Inches(0.5), Inches(3.3), Inches(7.5), Inches(3.7),
         code_lines, lang_hint="py",
         title="workers/katara_diagnostic/orchestrator.py", size=9)

# Encadré « pourquoi séquentiel »
add_rect(s, Inches(8.2), Inches(3.3), Inches(4.6), Inches(3.7),
         WHITE, line=GREY_LIGHT)
add_rect(s, Inches(8.2), Inches(3.3), Inches(0.12), Inches(3.7), GOLD)
add_text(s, Inches(8.4), Inches(3.4), Inches(4.3), Inches(0.4),
         "Pourquoi séquentiel ?", size=13, bold=True, color=NAVY)
add_bullets(s, Inches(8.4), Inches(3.85), Inches(4.3), Inches(3.0), [
    "Gemini = goulot (~5-15s)",
    "Gain d'un asyncio.gather < 100 ms",
    "Coût : raisonnement sur les modes d'échec ↑",
    "→ tradeoff lisibilité > micro-perf",
], size=11, color=GREY)

add_text(s, Inches(8.4), Inches(5.3), Inches(4.3), Inches(0.4),
         "Anti prompt-injection", size=13, bold=True, color=NAVY)
add_bullets(s, Inches(8.4), Inches(5.7), Inches(4.3), Inches(1.4), [
    "Jinja2 autoescape activé",
    "Templates .j2 dédiés (FR/AR/EN)",
    "Données utilisateur nettoyées",
], size=11, color=GREY)

footer(s, 14, "02 — Backend & Frontend")


# ============================================================
# Slide 15 — Frontend stack
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, LIGHT)
header_band(s, "Frontend — Next.js 15 / React 19",
            "App Router, Server Components, Tailwind v4", axis="AXE 02")

# 3 cartes haut : stack / archi / qualité
def card(slide, x, y, w, h, title, lines, accent):
    add_rect(slide, x, y, w, h, WHITE, line=GREY_LIGHT)
    add_rect(slide, x, y, Inches(0.1), h, accent)
    add_text(slide, x + Inches(0.25), y + Inches(0.12), w - Inches(0.3),
             Inches(0.4), title, size=14, bold=True, color=NAVY)
    add_bullets(slide, x + Inches(0.25), y + Inches(0.55),
                w - Inches(0.3), h - Inches(0.6), lines, size=11, color=GREY)


card(s, Inches(0.4), Inches(1.3), Inches(4.1), Inches(3.0),
     "Stack", [
        "Next.js 15.1 (App Router)",
        "React 19 — Server Components",
        "TypeScript 5.7 (strict)",
        "Tailwind v4 + PostCSS",
        "@supabase/ssr (cookies)",
        "Zod — validation formulaires",
        "Leaflet — cartes parcelles",
        "Sentry — erreurs client + edge",
     ], TEAL)

card(s, Inches(4.65), Inches(1.3), Inches(4.1), Inches(3.0),
     "Architecture", [
        "Middleware d'auth global (Edge)",
        "Dashboard segmenté par rôle :",
        "  /dashboard/farmer",
        "  /dashboard/restaurant",
        "  /dashboard/admin",
        "Server Actions + RPC Supabase",
        "Routes API minimales (/api/…)",
        "25 pages page.tsx au total",
     ], GOLD)

card(s, Inches(8.9), Inches(1.3), Inches(4.0), Inches(3.0),
     "Qualité & tests", [
        "Vitest + Testing Library",
        "161 fichiers de tests frontend",
        "ESLint (next/core-web-vitals)",
        "Type-check CI (tsc --noEmit)",
        "Sentry source-maps build",
        "Dockerfile multi-stage",
     ], NAVY)

# Écrans livrés
add_rect(s, Inches(0.4), Inches(4.4), Inches(12.5), Inches(2.55), WHITE,
         line=GREY_LIGHT)
add_rect(s, Inches(0.4), Inches(4.4), Inches(0.12), Inches(2.55), TEAL)
add_text(s, Inches(0.65), Inches(4.5), Inches(11.8), Inches(0.4),
         "Écrans livrés par rôle", size=14, bold=True, color=NAVY)

screens = [
    ("Auth (4)",
     "Login • Register • Onboarding • Vérification KYC (upload + statut)"),
    ("Farmer (6)",
     "Dashboard • Parcelles (carte Leaflet GeoJSON) • Devices (pair, rotate, unlink) "
     "• Télémétrie live • Diagnostics IA • Seuils d'alerte"),
    ("Restaurant (5)",
     "Catalogue FarMarket (filtres) • Détail produit • Panier • Pipeline commandes "
     "• Tracking livraison"),
    ("Admin (4)",
     "Modération KYC (queue) • Gestion FarMarket • Vue système • Utilisateurs"),
]
y = Inches(4.95)
for who, what in screens:
    add_rect(s, Inches(0.8), y, Inches(1.5), Inches(0.4), GOLD)
    add_text(s, Inches(0.8), y + Inches(0.05), Inches(1.5), Inches(0.3),
             who, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(2.5), y + Inches(0.05), Inches(10), Inches(0.4),
             what, size=10, color=GREY)
    y += Inches(0.46)
footer(s, 15, "02 — Backend & Frontend")


# ============================================================
# Slide 16 — PREUVE 5 — Middleware Next.js
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, LIGHT)
header_band(s, "Preuve 5 — Middleware Next.js",
            "Edge runtime — auth gate + vérification KYC par décodage JWT",
            axis="AXE 02")

code_lines = [
    "const PROTECTED_PREFIXES = ['/dashboard', '/onboarding/verification', '/admin'];",
    "const VERIFIED_PRO_PREFIXES = ['/farmarket/new', '/secondserve/new'];",
    "",
    "export async function middleware(request: NextRequest) {",
    "  let response = NextResponse.next({ request });",
    "  const supabase = createServerClient(",
    "    process.env.NEXT_PUBLIC_SUPABASE_URL!,",
    "    process.env.NEXT_PUBLIC_SUPABASE_ANON_KEY!,",
    "    { cookies: { /* ... refresh des tokens via setAll ... */ } },",
    "  );",
    "",
    "  // IMPORTANT : getUser() (et pas getSession) → revalide le cookie",
    "  const { data: { user } } = await supabase.auth.getUser();",
    "  const { pathname } = request.nextUrl;",
    "",
    "  // Route protégée + pas de user → /login?next=...",
    "  if (isProtected(pathname) && !user) {",
    "    const url = request.nextUrl.clone();",
    "    url.pathname = '/login';",
    "    url.searchParams.set('next', pathname);",
    "    return NextResponse.redirect(url);",
    "  }",
    "",
    "  // Route 'pro' (FarMarket/SecondServe) → exige verification VERIFIED",
    "  if (isPublishRoute(pathname) && user) {",
    "    const { data: { session } } = await supabase.auth.getSession();",
    "    const claims = decodeJwtClaims(session!.access_token);",
    "    const isPro = claims.user_role === 'FARMER' || claims.user_role === 'RESTAURANT';",
    "    if (isPro && claims.verification_status !== 'VERIFIED') {",
    "      return NextResponse.redirect('/onboarding/verification');",
    "    }",
    "  }",
    "  return response;",
    "}",
]
code_box(s, Inches(0.5), Inches(1.3), Inches(8.5), Inches(5.7),
         code_lines, lang_hint="ts",
         title="frontend/src/middleware.ts", size=9)

# Encadré droite
add_rect(s, Inches(9.2), Inches(1.3), Inches(3.6), Inches(2.7),
         WHITE, line=GREY_LIGHT)
add_rect(s, Inches(9.2), Inches(1.3), Inches(0.12), Inches(2.7), GREEN)
add_text(s, Inches(9.4), Inches(1.4), Inches(3.4), Inches(0.4),
         "Pourquoi le middleware ?", size=12, bold=True, color=NAVY)
add_bullets(s, Inches(9.4), Inches(1.8), Inches(3.4), Inches(2.3), [
    "Tourne à l'Edge (faible latence)",
    "Refresh cookies sur CHAQUE requête",
    "Redirige AVANT que la page ne render",
    "UX uniquement — pas la sécurité",
], size=10, color=GREY)

add_rect(s, Inches(9.2), Inches(4.2), Inches(3.6), Inches(2.8),
         WHITE, line=GREY_LIGHT)
add_rect(s, Inches(9.2), Inches(4.2), Inches(0.12), Inches(2.8), RED)
add_text(s, Inches(9.4), Inches(4.3), Inches(3.4), Inches(0.4),
         "Vraie sécurité : 2 couches", size=12, bold=True, color=NAVY)
add_bullets(s, Inches(9.4), Inches(4.7), Inches(3.4), Inches(2.3), [
    "RLS Postgres (read+write)",
    "Dépendance require_verified() FastAPI",
    "Le middleware est juste UX",
    "→ Si un user contourne le middleware,",
    "   la DB refuse de toute façon",
], size=10, color=GREY)
footer(s, 16, "02 — Backend & Frontend")


# ============================================================
# Slide 17 — Communication Front ↔ Back
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, LIGHT)
header_band(s, "Communication Front ↔ Back",
            "Deux canaux complémentaires — sécurisés par RLS Postgres",
            axis="AXE 02")

# Canal 1
add_rect(s, Inches(0.5), Inches(1.4), Inches(6.1), Inches(5.4),
         WHITE, line=GREY_LIGHT)
add_rect(s, Inches(0.5), Inches(1.4), Inches(6.1), Inches(0.5), TEAL)
add_text(s, Inches(0.5), Inches(1.5), Inches(6.1), Inches(0.4),
         "① Direct Supabase (CRUD simple)",
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_bullets(s, Inches(0.75), Inches(2.1), Inches(5.6), Inches(3.0), [
    "Frontend ↔ PostgREST via @supabase/ssr",
    "JWT propagé via cookies HttpOnly",
    "RLS appliquée côté base — pas de leak",
    "Realtime channels (futur)",
    "Storage SDK (upload KYC, photos)",
], size=11, color=GREY)
add_text(s, Inches(0.75), Inches(4.5), Inches(5.6), Inches(0.4),
         "Quand l'utiliser :", size=12, bold=True, color=NAVY)
add_bullets(s, Inches(0.75), Inches(4.85), Inches(5.6), Inches(1.8), [
    "Lectures de listes (ads, parcelles)",
    "Mutations triviales protégées par RLS",
    "Upload de fichiers vers Storage",
], size=11, color=GREEN)

# Canal 2
add_rect(s, Inches(6.7), Inches(1.4), Inches(6.1), Inches(5.4),
         WHITE, line=GREY_LIGHT)
add_rect(s, Inches(6.7), Inches(1.4), Inches(6.1), Inches(0.5), GOLD)
add_text(s, Inches(6.7), Inches(1.5), Inches(6.1), Inches(0.4),
         "② FastAPI (logique métier + intégrations)",
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_bullets(s, Inches(6.95), Inches(2.1), Inches(5.6), Inches(3.0), [
    "Endpoints REST documentés (OpenAPI auto)",
    "Vérif JWT côté serveur (PyJWT + JWKS)",
    "Orchestration : devices, diagnostics IA, KYC",
    "Intégrations : Gemini, Sentinel-2, Brevo, OWM",
    "Workers asynchrones via NOTIFY",
], size=11, color=GREY)
add_text(s, Inches(6.95), Inches(4.5), Inches(5.6), Inches(0.4),
         "Quand l'utiliser :", size=12, bold=True, color=NAVY)
add_bullets(s, Inches(6.95), Inches(4.85), Inches(5.6), Inches(1.8), [
    "Pairing device (génère + bcrypt côté serveur)",
    "Déclencher diagnostic IA (workflow)",
    "Modération KYC (transitions admin)",
    "Tout ce qui touche un secret ou un service externe",
], size=11, color=GREEN)
footer(s, 17, "02 — Backend & Frontend")


# ============================================================
# SECTION 03 — SÉCURITÉ
# ============================================================
s = section_divider("03", "Sécurité",
                    "Défense en profondeur — DB, API, Edge, devices", 3)
items = [
    "Les 4 piliers — une couche défensive par niveau du stack",
    "PREUVE 6 — Bcrypt + pgcrypto.crypt (constant-time)",
    "PREUVE 7 — RLS policy katara_devices_insert (vraie SQL)",
    "Surface d'attaque : 7 risques identifiés + leurs mitigations",
    "Tests & garde-fous (allowlist service_role, contract test)",
]
y = Inches(4.95)
for it in items:
    add_text(s, Inches(0.8), y, Inches(12), Inches(0.4),
             "›  " + it, size=14, color=WHITE)
    y += Inches(0.35)


# ============================================================
# Slide 19 — 4 piliers
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, LIGHT)
header_band(s, "Sécurité — 4 piliers", "Défense en profondeur",
            axis="AXE 03")

pillars = [
    ("AUTH", "Authentification",
     ["Supabase Auth (email/password)",
      "JWT HS256 signé",
      "Cookies HttpOnly + Secure",
      "Refresh token rotation",
      "Blocage signup admin (trigger SQL)",
      "Headers : WWW-Authenticate Bearer"], TEAL),
    ("AUTHZ", "Autorisation",
     ["JWT custom claim 'user_role'",
      "Auth Hook PL/pgSQL (mig. 0006)",
      "Fonction SQL has_role()",
      "RLS sur 100% des tables (54 policies)",
      "FORCE ROW LEVEL SECURITY",
      "Vues SECURITY INVOKER"], GOLD),
    ("DATA", "Données & Storage",
     ["Buckets séparés (privé / public)",
      "Storage policies RLS",
      "Audit guards SQL (RAISE)",
      "Migration bookkeeping RLS",
      "service_role allowlist testée",
      "Verification claim AUTH-06"], NAVY),
    ("DEVICES", "Devices IoT",
     ["Clé vk_<32 hex> (128 bits)",
      "Hash bcrypt cost=10",
      "Affichée 1x au pairing",
      "pgcrypto.crypt constant-time",
      "Endpoint ingest = service_role",
      "Rotate + unlink supportés"], RED),
]
xs = [Inches(0.4), Inches(3.65), Inches(6.9), Inches(10.15)]
for (code, name, items_, color), x in zip(pillars, xs):
    add_rect(s, x, Inches(1.3), Inches(3.05), Inches(5.7), WHITE,
             line=GREY_LIGHT)
    add_rect(s, x, Inches(1.3), Inches(3.05), Inches(0.7), color)
    add_text(s, x, Inches(1.4), Inches(3.05), Inches(0.5),
             code, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.15), Inches(2.15), Inches(2.85), Inches(0.5),
             name, size=14, bold=True, color=NAVY)
    add_bullets(s, x + Inches(0.15), Inches(2.6), Inches(2.85), Inches(4.2),
                items_, size=10, color=GREY)
footer(s, 19, "03 — Sécurité")


# ============================================================
# Slide 20 — PREUVE 6 — Bcrypt + pgcrypto constant-time
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, LIGHT)
header_band(s, "Preuve 6 — Bcrypt + pgcrypto (constant-time)",
            "La clé device n'est JAMAIS stockée en clair", axis="AXE 03")

# Python à gauche
py_code = [
    "# backend/app/core/api_keys.py",
    "_PREFIX = 'vk_'",
    "_HEX_BYTES = 16          # 128 bits d'entropie",
    "_BCRYPT_COST = 10        # ~10 ms sur VPS demo",
    "",
    "def generate_device_api_key() -> str:",
    "    return _PREFIX + secrets.token_hex(_HEX_BYTES)",
    "",
    "def hash_device_api_key(plaintext: str) -> str:",
    "    raw = bcrypt.hashpw(",
    "        plaintext.encode('utf-8'),",
    "        bcrypt.gensalt(rounds=_BCRYPT_COST),",
    "    )",
    "    # Python bcrypt émet $2b$, pgcrypto attend $2a$",
    "    # Mêmes algos, juste prefix différent",
    "    return raw.decode().replace('$2b$', '$2a$', 1)",
    "",
    "def last4(plaintext: str) -> str:",
    "    return plaintext[-4:]    # UI display only",
]
code_box(s, Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.7),
         py_code, lang_hint="py",
         title="① Génération + hash côté Python (FastAPI)", size=10)

# SQL à droite
sql_code = [
    "create function public.verify_device_api_key(",
    "    p_device_id text,",
    "    p_api_key   text",
    ") returns table (",
    "    device_row_id uuid,",
    "    parcel_id     uuid,",
    "    farmer_id     uuid",
    ")",
    "language sql stable security definer",
    "set search_path = public, pg_temp",
    "as $$",
    "    select d.id, d.parcel_id, d.farmer_id",
    "      from public.m1_katara_devices d",
    "     where d.device_id    = p_device_id",
    "       and d.status      <> 'UNLINKED'",
    "       -- crypt() = constant-time compare",
    "       -- en C dans pgcrypto",
    "       and d.api_key_hash =",
    "           extensions.crypt(p_api_key, d.api_key_hash)",
    "     limit 1;",
    "$$;",
    "",
    "revoke all on function public.verify_device_api_key(text, text)",
    "    from public;",
    "grant execute on function public.verify_device_api_key(text, text)",
    "    to service_role;",
]
code_box(s, Inches(6.7), Inches(1.3), Inches(6.1), Inches(5.7),
         sql_code, lang_hint="sql",
         title="② Vérification côté Postgres (migration 0017)", size=10)
footer(s, 20, "03 — Sécurité")


# ============================================================
# Slide 21 — PREUVE 7 — RLS policy
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, LIGHT)
header_band(s, "Preuve 7 — RLS policy (vraie SQL en prod)",
            "katara_devices_insert — 3 conditions cumulatives, contrat AUTH-04",
            axis="AXE 03")

# Code SQL
sql_code = [
    "alter table public.m1_katara_devices",
    "    enable row level security;",
    "    -- + FORCE (via trigger AUTH-04, slide 8)",
    "",
    "create policy katara_devices_insert_verified_farmer_owns_parcel",
    "    on public.m1_katara_devices",
    "    for insert to authenticated",
    "    with check (",
    "        -- ① le caller est bien le farmer revendiqué",
    "        auth.uid() = farmer_id",
    "",
    "        -- ② et il a le rôle FARMER",
    "        and public.has_role('FARMER'::public.user_role)",
    "",
    "        -- ③ et son KYC est VERIFIED (claim AUTH-06)",
    "        and (",
    "            select verification_status",
    "              from public.profiles",
    "             where id = auth.uid()",
    "        ) = 'VERIFIED'",
    "",
    "        -- ④ et la parcelle visée lui appartient bien",
    "        and exists (",
    "            select 1",
    "              from public.m1_katara_parcels p",
    "             where p.id = parcel_id",
    "               and p.farmer_id = auth.uid()",
    "        )",
    "    );",
]
code_box(s, Inches(0.5), Inches(1.3), Inches(8.0), Inches(5.7),
         sql_code, lang_hint="sql",
         title="db/migrations/0017_kat02_katara_devices.sql", size=10)

# Explications
add_rect(s, Inches(8.7), Inches(1.3), Inches(4.1), Inches(5.7),
         WHITE, line=GREY_LIGHT)
add_rect(s, Inches(8.7), Inches(1.3), Inches(0.12), Inches(5.7), TEAL)
add_text(s, Inches(8.95), Inches(1.4), Inches(3.7), Inches(0.4),
         "Lecture ligne par ligne", size=13, bold=True, color=NAVY)
add_bullets(s, Inches(8.95), Inches(1.85), Inches(3.7), Inches(2.8), [
    "① Anti-spoofing du farmer_id",
    "② Rôle (du JWT claim)",
    "③ KYC réellement validé",
    "④ Possession effective de la parcelle",
], size=11, color=GREY)

add_text(s, Inches(8.95), Inches(4.8), Inches(3.7), Inches(0.4),
         "Défense en profondeur", size=13, bold=True, color=NAVY)
add_bullets(s, Inches(8.95), Inches(5.2), Inches(3.7), Inches(1.7), [
    "Le code Python valide les mêmes conditions",
    "→ deux barrières indépendantes",
    "→ une couche peut tomber sans leak",
], size=11, color=GREY)
footer(s, 21, "03 — Sécurité")


# ============================================================
# Slide 22 — Surface d'attaque
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, LIGHT)
header_band(s, "Surface d'attaque & mitigations",
            "Risques identifiés → contre-mesures en place",
            axis="AXE 03")

risks = [
    ("Vol de clé device ESP32",
     "Hash bcrypt cost=10 en base, plaintext jamais persisté. Ingest = service_role uniquement. Unlink + rotate."),
    ("Injection prompt Gemini",
     "Jinja2 autoescape=on sur tous templates. Inputs utilisateur tronqués, templates dédiés FR/AR/EN."),
    ("Élévation privilèges via signup",
     "Trigger SQL AUTH-02 bloque l'auto-signup avec role=admin. Création admin par script CLI."),
    ("Leak documents KYC",
     "Bucket privé kyc-private + Storage policies RLS. Signed URLs courte durée pour preview admin."),
    ("Bypass RLS par owner BDD",
     "FORCE ROW LEVEL SECURITY → personne n'échappe, même superuser. service_role usage allowlisté."),
    ("Énumération devices via timing",
     "Erreur constante _INVALID_CREDS, bcrypt verify constant-time (pgcrypto en C), pas de leak d'existence."),
    ("XSS / CSRF frontend",
     "React 19 escape par défaut, Zod valide tous formulaires, cookies HttpOnly + SameSite=Lax, CORS strict."),
]
y = Inches(1.3)
for risk, mitig in risks:
    add_rect(s, Inches(0.4), y, Inches(0.4), Inches(0.75), RED)
    add_text(s, Inches(0.4), y + Inches(0.2), Inches(0.4), Inches(0.4),
             "!", size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.95), y + Inches(0.05), Inches(5), Inches(0.4),
             risk, size=12, bold=True, color=NAVY)
    add_text(s, Inches(0.95), y + Inches(0.4), Inches(11.8), Inches(0.4),
             mitig, size=10, color=GREY)
    y += Inches(0.82)
footer(s, 22, "03 — Sécurité")


# ============================================================
# Slide 23 — Tests & garde-fous
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, LIGHT)
header_band(s, "Tests & garde-fous de sécurité",
            "Ce qui empêche une régression de passer en prod", axis="AXE 03")

cols = [
    ("Tests Postgres (pgTAP-like)", TEAL, [
        "auth04_rls_contract.sql : assertions sur 100% tables",
        "kat02_devices_rls.sql : 16 scénarios cross-farmer",
        "kat03_ingest_force_rls.sql : INSERT direct refusé",
        "Rejoués à chaque CI sur DB de test fraîche",
    ]),
    ("Tests Python (pytest)", GOLD, [
        "41 fichiers test_*.py côté backend",
        "test_service_client_callsite_allowlist.py :",
        "  → grep tous les usages de service_role",
        "  → fail si callsite non documenté",
        "test_security_ingest.py : énumération impossible",
    ]),
    ("Garde-fous CI/CD", NAVY, [
        "scripts/verify-rls-enabled.sh (shell)",
        "Lint : ruff + mypy strict",
        "Type-check : tsc --noEmit côté Next.js",
        "Sentry alerte sur erreurs prod",
        "Pre-commit hooks (format + lint + tests rapides)",
    ]),
]
xs = [Inches(0.4), Inches(4.65), Inches(8.9)]
for (title, color, items_), x in zip(cols, xs):
    add_rect(s, x, Inches(1.3), Inches(4.1), Inches(5.7), WHITE,
             line=GREY_LIGHT)
    add_rect(s, x, Inches(1.3), Inches(4.1), Inches(0.6), color)
    add_text(s, x, Inches(1.45), Inches(4.1), Inches(0.5),
             title, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_bullets(s, x + Inches(0.2), Inches(2.1), Inches(3.7), Inches(4.5),
                items_, size=11, color=GREY)

# Bandeau bas
add_rect(s, Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.2), GREEN)
footer(s, 23, "03 — Sécurité")


# ============================================================
# Slide 24 — Roadmap
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, LIGHT)
header_band(s, "Prochaines étapes",
            "Roadmap restante avant soutenance + post-PFE")

cols = [
    ("Court terme  (2-3 sem.)", TEAL, [
        "KAT — stories pH/EC restantes (schema update)",
        "FAR-11 → FAR-13 : factures PDF, litiges",
        "CI/CD : déploiement auto Vercel + Fly.io",
        "Tests E2E Playwright sur parcours critiques",
        "Documentation API publique",
    ]),
    ("Moyen terme  (1-2 mois)", GOLD, [
        "Module SecondServe (invendus restau.)",
        "Application mobile (PWA installable)",
        "Notifications push (Web Push API)",
        "Dashboard analytics admin (charts)",
        "Internationalisation AR + EN complète",
    ]),
    ("Industrialisation", NAVY, [
        "Audit sécurité externe (pen-test)",
        "Monitoring : Grafana + Prometheus",
        "Backups Postgres auto + tests restore",
        "Documentation technique complète",
        "ESP32 OTA (firmware update à distance)",
    ]),
]
xs = [Inches(0.4), Inches(4.65), Inches(8.9)]
for (title, color, items_), x in zip(cols, xs):
    add_rect(s, x, Inches(1.3), Inches(4.1), Inches(5.7), WHITE,
             line=GREY_LIGHT)
    add_rect(s, x, Inches(1.3), Inches(4.1), Inches(0.7), color)
    add_text(s, x, Inches(1.45), Inches(4.1), Inches(0.5),
             title, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_bullets(s, x + Inches(0.2), Inches(2.2), Inches(3.7), Inches(4.5),
                items_, size=12, color=GREY)
footer(s, 24, "Roadmap")


# ============================================================
# Slide 25 — Merci / Q&A
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, NAVY)
add_rect(s, 0, Inches(4.2), SW, Inches(0.08), TEAL)
add_text(s, Inches(0.8), Inches(2.3), Inches(12), Inches(1.2),
         "Merci pour votre attention", size=52, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(3.3), Inches(12), Inches(0.6),
         "Questions & démonstration live", size=22, color=TEAL)
add_text(s, Inches(0.8), Inches(4.5), Inches(12), Inches(0.4),
         "Yasser  —  PFE VitaChain  —  Mai 2026",
         size=15, color=GOLD)
add_text(s, Inches(0.8), Inches(5.0), Inches(12), Inches(0.4),
         "yasseralgoside@gmail.com", size=12, color=WHITE)

# Pense-bête démo
add_rect(s, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.1),
         CODE_BG)
add_text(s, Inches(1.0), Inches(5.9), Inches(11.5), Inches(0.4),
         "Démo live au programme :", size=11, bold=True, color=GOLD,
         font="Consolas", italic=True)
add_text(s, Inches(1.0), Inches(6.25), Inches(11.5), Inches(0.6),
         "1. Pair ESP32 → 2. Voir télémétrie temps réel → 3. Diagnostic IA → "
         "4. Cycle de commande FarMarket complet",
         size=11, color=WHITE, font="Consolas")


# ============================================================
out = Path(__file__).parent / "VitaChain_Avancement.pptx"
prs.save(out)
print(f"OK — {out}  ({len(prs.slides)} slides)")
