"""
VitaChain — Générateur de présentation PowerPoint
Thème : blanc, minimaliste, accents bleu/vert VitaChain
Exécuter : python docs/generate_pptx.py
Sortie   : docs/presentation_vitachain.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── Palette VitaChain ────────────────────────────────────────
BLUE   = RGBColor(37,  99,  235)   # vitaBlue
GREEN  = RGBColor(22,  163, 74)    # vitaGreen
DARK   = RGBColor(30,  58,  138)   # vitaDark
GRAY   = RGBColor(107, 114, 128)   # texte secondaire
LIGHT  = RGBColor(239, 246, 255)   # fond blocs
WHITE  = RGBColor(255, 255, 255)
BLACK  = RGBColor(15,  23,  42)

# ── Dimensions 16:9 ─────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # layout vide


# ─────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)          # MSO_SHAPE_TYPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, size=18, bold=False,
             color=BLACK, align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tf_box = slide.shapes.add_textbox(x, y, w, h)
    tf = tf_box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tf_box


def add_bullet_block(slide, items, x, y, w, h,
                     size=14, color=BLACK, dot_color=BLUE, spacing=Pt(6)):
    tf_box = slide.shapes.add_textbox(x, y, w, h)
    tf = tf_box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = spacing
        run_dot = p.add_run()
        run_dot.text = "● "
        run_dot.font.size = Pt(size - 2)
        run_dot.font.color.rgb = dot_color
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tf_box


def slide_header(slide, title, subtitle=None):
    """Barre d'en-tête bleue + filet vert + titre."""
    add_rect(slide, 0, 0, W, Inches(1.22), BLUE)
    add_rect(slide, 0, Inches(1.22), W, Inches(0.06), GREEN)
    add_text(slide, title,
             Inches(0.4), Inches(0.18), W - Inches(0.8), Inches(0.75),
             size=26, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.4), Inches(0.76), W - Inches(0.8), Inches(0.4),
                 size=13, color=RGBColor(186, 210, 255), align=PP_ALIGN.LEFT)


def slide_footer(slide):
    """Filet bas + numéro de slide."""
    add_rect(slide, 0, H - Inches(0.35), W, Inches(0.04), GREEN)
    add_rect(slide, 0, H - Inches(0.35) + Inches(0.04), W, Inches(0.31), BLUE)
    add_text(slide,
             "Badre Saad  |  Kodoussi Mohammed  |  El Karmi Yasser          "
             "VitaChain — PFE 2026          Faculté des Sciences Ben M'Sik",
             Inches(0.3), H - Inches(0.32), W - Inches(0.6), Inches(0.28),
             size=8, color=RGBColor(186, 210, 255), align=PP_ALIGN.CENTER)


def info_card(slide, title, body, x, y, w, h,
              bg=LIGHT, border=BLUE, title_color=BLUE, body_color=BLACK):
    """Carte avec titre coloré + corps."""
    add_rect(slide, x, y, w, Inches(0.03), border)
    bg_shape = add_rect(slide, x, y + Inches(0.03), w, h - Inches(0.03), bg)
    add_text(slide, title,
             x + Inches(0.15), y + Inches(0.08),
             w - Inches(0.3), Inches(0.32),
             size=12, bold=True, color=title_color)
    add_text(slide, body,
             x + Inches(0.15), y + Inches(0.38),
             w - Inches(0.3), h - Inches(0.5),
             size=11, color=body_color)


def kpi_badge(slide, value, label, x, y, w=Inches(1.9), h=Inches(0.9),
              color=BLUE):
    add_rect(slide, x, y, w, h, LIGHT)
    add_rect(slide, x, y, w, Inches(0.04), color)
    add_text(slide, value,
             x, y + Inches(0.04), w, Inches(0.5),
             size=22, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, label,
             x, y + Inches(0.52), w, Inches(0.35),
             size=9, color=GRAY, align=PP_ALIGN.CENTER)


# =============================================================
# SLIDE 1 — PAGE DE TITRE
# =============================================================
sl = prs.slides.add_slide(BLANK)

# Fond blanc implicite
# Bande bleue supérieure (40%)
add_rect(sl, 0, 0, W, Inches(3.0), BLUE)
add_rect(sl, 0, Inches(3.0), W, Inches(0.08), GREEN)

# Nom de la plateforme
add_text(sl, "VitaChain",
         Inches(0.6), Inches(0.4), Inches(10), Inches(1.1),
         size=54, bold=True, color=WHITE)

add_text(sl, "Plateforme numérique intégrée pour l'agriculture intelligente\n"
             "et la chaîne d'approvisionnement alimentaire",
         Inches(0.6), Inches(1.5), Inches(10), Inches(1.2),
         size=18, color=RGBColor(186, 210, 255))

# Pastilles modules
modules = [
    ("Katara", "IoT Agricole",      BLUE,                  Inches(0.6)),
    ("FarMarket", "Marketplace B2B", RGBColor(22,163,74),  Inches(3.2)),
    ("SecondServe","Anti-gaspillage",RGBColor(234,88,12),  Inches(5.8)),
    ("Administration","KYC & Admin", GRAY,                 Inches(8.4)),
]
for name, sub, col, xpos in modules:
    add_rect(sl, xpos, Inches(2.3), Inches(2.3), Inches(0.55), col)
    add_text(sl, f"{name}  —  {sub}",
             xpos + Inches(0.12), Inches(2.33),
             Inches(2.1), Inches(0.5),
             size=10, bold=True, color=WHITE)

# Bloc inférieur
add_text(sl, "Réalisé par :   Badre Saad  ·  Kodoussi Mohammed  ·  El Karmi Yasser",
         Inches(0.6), Inches(3.55), Inches(7.5), Inches(0.45),
         size=13, bold=True, color=BLACK)
add_text(sl, "Encadré par :   Mme Ouchra Hafssa  ·  Mme Achtaich Khadija  ·  M. Ait Daoud Mohammed",
         Inches(0.6), Inches(3.98), Inches(11), Inches(0.4),
         size=12, color=GRAY)
add_text(sl, "Faculté des Sciences Ben M'Sik  —  Département Informatique  —  DIARS  —  2025 / 2026",
         Inches(0.6), Inches(4.38), Inches(11), Inches(0.4),
         size=11, italic=True, color=GRAY)
add_text(sl, "Soutenu le 19 juin 2026",
         Inches(0.6), Inches(4.80), Inches(5), Inches(0.38),
         size=11, color=GRAY)

# Filet bas
add_rect(sl, 0, H - Inches(0.25), W, Inches(0.05), GREEN)
add_rect(sl, 0, H - Inches(0.20), W, Inches(0.20), BLUE)


# =============================================================
# SLIDE 2 — PROBLÉMATIQUE
# =============================================================
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Problématique & Contexte",
             "Trois défis structurels dans le secteur agroalimentaire marocain")

COLS = [Inches(0.4), Inches(4.45), Inches(8.5)]
TITLES = [
    "Manque de données\nagronomiques",
    "Désintermédiation\ninsuffisante",
    "Gaspillage\nalimentaire",
]
BODIES = [
    "Les agriculteurs décident sans mesures précises du sol : humidité, pH, température, conductivité.",
    "Des intermédiaires captent la valeur. Peu de contact direct entre agriculteurs et restaurateurs.",
    "Les invendus des restaurants partent au rebut faute d'un canal de redistribution adapté.",
]
COLS_C = [BLUE, GREEN, RGBColor(234, 88, 12)]

for i in range(3):
    add_rect(sl, COLS[i], Inches(1.45), Inches(3.7), Inches(3.3), COLS_C[i])
    add_rect(sl, COLS[i], Inches(1.45) + Inches(0.28), Inches(3.7), Inches(3.02), WHITE)
    add_text(sl, TITLES[i],
             COLS[i] + Inches(0.15), Inches(1.55),
             Inches(3.4), Inches(0.65),
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, BODIES[i],
             COLS[i] + Inches(0.15), Inches(1.45) + Inches(0.4),
             Inches(3.4), Inches(2.4),
             size=13, color=BLACK)

# Flèche solution
add_rect(sl, Inches(4.9), Inches(5.05), Inches(3.5), Inches(0.05), BLUE)
add_text(sl, "Solution : VitaChain — un seul écosystème, trois réponses complémentaires",
         Inches(0.4), Inches(5.25), W - Inches(0.8), Inches(0.5),
         size=14, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# Pastilles
mods = [
    ("Katara", BLUE,                Inches(1.4)),
    ("FarMarket", GREEN,            Inches(4.7)),
    ("SecondServe", RGBColor(234,88,12), Inches(7.8)),
    ("Administration", GRAY,        Inches(10.5)),
]
for name, col, xp in mods:
    add_rect(sl, xp, Inches(5.9), Inches(2.2), Inches(0.5), col)
    add_text(sl, name, xp + Inches(0.1), Inches(5.92),
             Inches(2.0), Inches(0.46),
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

slide_footer(sl)


# =============================================================
# SLIDE 3 — ARCHITECTURE & INFRASTRUCTURE
# =============================================================
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Architecture & Infrastructure",
             "Déploiement production · Docker · NGINX · TLS · CI/CD")

# Diagramme N-tiers (gauche)
tiers = [
    ("Tier 1  —  Présentation",  "Next.js 15  +  SecondServe (Vite 6)",  BLUE,                Inches(1.45)),
    ("Tier 2  —  API",            "FastAPI (Python 3.12)  +  Workers async", GREEN,             Inches(2.35)),
    ("Tier 3  —  Données",        "PostgreSQL 17  +  Supabase Auth & Storage", RGBColor(109,40,217), Inches(3.25)),
    ("Tier 4  —  Infrastructure", "Docker · NGINX 1.27 · Let's Encrypt TLS", GRAY,             Inches(4.15)),
]
for label, sub, col, ypos in tiers:
    add_rect(sl, Inches(0.4), ypos, Inches(5.8), Inches(0.72), col)
    add_rect(sl, Inches(0.4), ypos + Inches(0.06), Inches(5.8), Inches(0.66), LIGHT)
    add_rect(sl, Inches(0.4), ypos, Inches(0.06), Inches(0.72), col)
    add_text(sl, label,
             Inches(0.6), ypos + Inches(0.04),
             Inches(2.8), Inches(0.3),
             size=10, bold=True, color=col)
    add_text(sl, sub,
             Inches(0.6), ypos + Inches(0.34),
             Inches(5.4), Inches(0.3),
             size=10, color=GRAY)

# Flèches entre tiers
for y in [Inches(2.17), Inches(3.07), Inches(3.97)]:
    add_rect(sl, Inches(3.15), y, Inches(0.06), Inches(0.18), GRAY)

# Points clés (droite)
add_text(sl, "Points clés",
         Inches(7.0), Inches(1.4), Inches(5.8), Inches(0.4),
         size=15, bold=True, color=DARK)

points = [
    ("HTTPS / TLS 1.3",     "NGINX + Certbot + Let's Encrypt — toutes les comm. chiffrées"),
    ("Docker Compose",       "6 services conteneurisés, orchestration complète"),
    ("CI/CD GitHub Actions", "5 jobs à chaque push : frontend, backend, DB, infra, secrets"),
    ("Sauvegardes nightly",  "pg_dump → rclone → Backblaze B2 chaque nuit"),
    ("Monitoring",           "Uptime Kuma (disponibilité) + Sentry (erreurs applicatives)"),
    ("Rate-limiting",        "NGINX : 30 req/s sur l'endpoint IoT"),
]
y0 = Inches(1.85)
for title, body in points:
    add_rect(sl, Inches(7.0), y0, Inches(0.06), Inches(0.6), BLUE)
    add_text(sl, title, Inches(7.2), y0, Inches(5.5), Inches(0.28),
             size=11, bold=True, color=BLUE)
    add_text(sl, body, Inches(7.2), y0 + Inches(0.26), Inches(5.5), Inches(0.28),
             size=10, color=GRAY)
    y0 += Inches(0.68)

slide_footer(sl)


# =============================================================
# SLIDE 4 — SÉCURITÉ : JWT & RLS
# =============================================================
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Sécurité : JWT multi-rôle, KYC & RLS",
             "Isolation des données garantie au niveau du moteur PostgreSQL")

# Colonne gauche : Auth & KYC
add_text(sl, "Authentification & KYC",
         Inches(0.4), Inches(1.42), Inches(5.8), Inches(0.38),
         size=14, bold=True, color=DARK)

auth_items = [
    "4 rôles : FARMER · RESTAURANT · CITIZEN · ADMIN",
    "JWT enrichi via Custom Access Token Hook Supabase\n(claims : user_role + verification_status)",
    "Durée de vie : 1 heure + rafraîchissement automatique",
    "Clés API IoT hachées bcrypt — jamais stockées en clair",
]
add_bullet_block(sl, auth_items, Inches(0.4), Inches(1.85),
                 Inches(5.8), Inches(2.0), size=12)

add_text(sl, "Processus KYC",
         Inches(0.4), Inches(3.95), Inches(5.8), Inches(0.38),
         size=14, bold=True, color=DARK)

kyc_steps = [
    "1.  Soumission documents (CIN / RC / carte agricole)",
    "2.  Upload sécurisé → Supabase Storage (RLS)",
    "3.  Revue par l'administrateur",
    "4.  Notification email via Brevo",
    "5.  Déblocage des actions professionnelles",
]
add_bullet_block(sl, kyc_steps, Inches(0.4), Inches(4.38),
                 Inches(5.8), Inches(2.1), size=12, dot_color=GREEN)

# Colonne droite : Matrice RLS
add_text(sl, "Matrice RLS — 22 politiques PostgreSQL",
         Inches(6.6), Inches(1.42), Inches(6.3), Inches(0.38),
         size=14, bold=True, color=DARK)

# En-têtes tableau
headers = ["Table", "FARMER", "RESTO", "CITIZEN", "ADMIN"]
col_xs  = [Inches(6.6), Inches(8.5), Inches(9.6), Inches(10.7), Inches(11.8)]
col_ws  = [Inches(1.85), Inches(1.0), Inches(1.0), Inches(1.0), Inches(1.0)]

add_rect(sl, Inches(6.6), Inches(1.85), Inches(6.3), Inches(0.38), BLUE)
for i, h in enumerate(headers):
    add_text(sl, h, col_xs[i], Inches(1.87), col_ws[i], Inches(0.34),
             size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rows = [
    ("profiles",    "own", "own", "own", "✓"),
    ("kyc_docs",    "own", "own", "✗",   "✓"),
    ("parcelles",   "own", "✗",   "✗",   "✓"),
    ("télémétrie",  "own", "✗",   "✗",   "✓"),
    ("annonces",    "✓",   "✓",   "✓",   "✓"),
    ("commandes",   "✗",   "own", "✗",   "✓"),
    ("items",       "anon","own", "✗",   "✓"),
    ("évaluations", "✓",   "D",   "✓",   "✓"),
]
for ri, row in enumerate(rows):
    bg = LIGHT if ri % 2 == 0 else WHITE
    add_rect(sl, Inches(6.6), Inches(2.23) + ri * Inches(0.46),
             Inches(6.3), Inches(0.46), bg)
    add_text(sl, row[0],
             col_xs[0], Inches(2.27) + ri * Inches(0.46),
             col_ws[0], Inches(0.38),
             size=10, color=BLACK)
    for ci, val in enumerate(row[1:], 1):
        col = GREEN if val == "✓" else (RGBColor(220,38,38) if val == "✗" else BLUE)
        add_text(sl, val,
                 col_xs[ci], Inches(2.27) + ri * Inches(0.46),
                 col_ws[ci], Inches(0.38),
                 size=10, bold=(val in ("✓","✗")), color=col,
                 align=PP_ALIGN.CENTER)

add_text(sl, "own = ses données uniquement  |  anon. = vue anonymisée sha256  |  D = achat livré requis",
         Inches(6.6), Inches(5.95), Inches(6.3), Inches(0.35),
         size=9, italic=True, color=GRAY)

slide_footer(sl)


# =============================================================
# SLIDE 5 — MODULE KATARA (IoT)
# =============================================================
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Module Katara — Surveillance IoT en temps réel",
             "ESP32 · Télémétrie sol · Alertes · Diagnostic IA · NDVI Satellite")

# Gauche : capteurs + pipeline
add_text(sl, "Dispositif ESP32 — capteurs intégrés",
         Inches(0.4), Inches(1.42), Inches(5.8), Inches(0.38),
         size=13, bold=True, color=DARK)
capteurs = [
    "Humidité du sol (capteur capacitif)",
    "Température du sol (DS18B20)",
    "pH du sol (électrode analogique)",
    "Conductivité électrique (EC meter)",
    "Niveau de batterie (ADC)",
]
add_bullet_block(sl, capteurs, Inches(0.4), Inches(1.85),
                 Inches(5.8), Inches(1.85), size=12)

# Pipeline
add_text(sl, "Pipeline d'ingestion",
         Inches(0.4), Inches(3.75), Inches(5.8), Inches(0.35),
         size=13, bold=True, color=DARK)
pipeline = [
    ("ESP32", BLUE),
    ("FastAPI", GREEN),
    ("PostgreSQL", RGBColor(109,40,217)),
    ("Dashboard", RGBColor(234,88,12)),
]
px = Inches(0.4)
for i, (name, col) in enumerate(pipeline):
    add_rect(sl, px, Inches(4.18), Inches(1.25), Inches(0.48), col)
    add_text(sl, name, px, Inches(4.20), Inches(1.25), Inches(0.44),
             size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 3:
        add_text(sl, "→", px + Inches(1.25), Inches(4.18),
                 Inches(0.2), Inches(0.48),
                 size=14, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    px += Inches(1.45)

add_text(sl, "HTTPS",
         Inches(1.65), Inches(4.7), Inches(1.0), Inches(0.3),
         size=8, italic=True, color=GRAY, align=PP_ALIGN.CENTER)
add_text(sl, "SQL",
         Inches(3.1), Inches(4.7), Inches(0.8), Inches(0.3),
         size=8, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

# Badge performance
add_rect(sl, Inches(0.4), Inches(5.2), Inches(5.8), Inches(0.6), BLUE)
add_text(sl, "  Latence d'ingestion : < 50 ms (p50)  —  fonction SQL atomique SECURITY DEFINER",
         Inches(0.4), Inches(5.24), Inches(5.8), Inches(0.52),
         size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Droite : fonctionnalités
add_text(sl, "Fonctionnalités avancées",
         Inches(6.6), Inches(1.42), Inches(6.3), Inches(0.38),
         size=13, bold=True, color=DARK)

cards = [
    ("Alertes automatiques",
     "Seuils min/max configurables par métrique → LISTEN/NOTIFY PostgreSQL → email Brevo",
     BLUE),
    ("Diagnostic IA — Google Gemini",
     "Télémétrie + météo OWM + NDVI Sentinel-2 → rapport agronomique en langage naturel",
     GREEN),
    ("Détection hors-ligne",
     "Worker CRON : capteur silencieux depuis > 1 h → alerte email automatique",
     RGBColor(234, 88, 12)),
    ("Sécurité IoT",
     "Clé API unique par dispositif (vk_<32hex>), hachée bcrypt, affichée une seule fois",
     RGBColor(109, 40, 217)),
]
yc = Inches(1.88)
for title, body, col in cards:
    add_rect(sl, Inches(6.6), yc, Inches(6.3), Inches(0.04), col)
    add_rect(sl, Inches(6.6), yc + Inches(0.04), Inches(6.3), Inches(1.0), LIGHT)
    add_text(sl, title, Inches(6.75), yc + Inches(0.08),
             Inches(6.0), Inches(0.3),
             size=11, bold=True, color=col)
    add_text(sl, body, Inches(6.75), yc + Inches(0.36),
             Inches(6.0), Inches(0.6),
             size=10, color=BLACK)
    yc += Inches(1.1)

slide_footer(sl)


# =============================================================
# SLIDE 6 — MODULE FARMARKET
# =============================================================
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Module FarMarket — Marketplace Agricole B2B",
             "Annonces · Catalogue · Commandes multi-agriculteurs · Anonymisation · Évaluations")

# Gauche : agriculteur
add_rect(sl, Inches(0.4), Inches(1.42), Inches(5.8), Inches(0.38), BLUE)
add_text(sl, "  Côté Agriculteur (FARMER vérifié)",
         Inches(0.4), Inches(1.44), Inches(5.8), Inches(0.34),
         size=12, bold=True, color=WHITE)
farmer = [
    "Publication d'annonces avec jusqu'à 5 photos",
    "Prix en MAD/kg · 12 régions du Maroc",
    "Expiration automatique à J+7",
    "Commandes entrantes anonymisées (resto_handle sha256)",
    "Acceptation / rejet par article",
]
add_bullet_block(sl, farmer, Inches(0.4), Inches(1.86),
                 Inches(5.8), Inches(1.95), size=12)

# Restaurateur
add_rect(sl, Inches(0.4), Inches(3.9), Inches(5.8), Inches(0.38), GREEN)
add_text(sl, "  Côté Restaurateur (RESTAURANT vérifié)",
         Inches(0.4), Inches(3.92), Inches(5.8), Inches(0.34),
         size=12, bold=True, color=WHITE)
resto = [
    "Catalogue paginé · filtres région / type / prix",
    "Recherche textuelle floue (extension pg_trgm)",
    "Panier multi-agriculteurs",
    "Paiement COD (à la livraison)",
    "Évaluation post-livraison : 1 à 5 étoiles",
]
add_bullet_block(sl, resto, Inches(0.4), Inches(4.34),
                 Inches(5.8), Inches(1.95), size=12, dot_color=GREEN)

# Droite : anonymisation + cycle de vie
add_rect(sl, Inches(6.6), Inches(1.42), Inches(6.3), Inches(0.04), RGBColor(234,88,12))
add_rect(sl, Inches(6.6), Inches(1.46), Inches(6.3), Inches(1.2), LIGHT)
add_text(sl, "Anonymisation B2B",
         Inches(6.75), Inches(1.5), Inches(6.0), Inches(0.32),
         size=12, bold=True, color=RGBColor(234,88,12))
add_text(sl, "Le restaurateur reste anonyme pour l'agriculteur via une vue SQL sécurisée :\n"
             "resto_handle = sha256(restaurant_id || farmer_id)",
         Inches(6.75), Inches(1.84), Inches(6.0), Inches(0.72),
         size=11, color=BLACK)

add_text(sl, "Cycle de vie d'une commande",
         Inches(6.6), Inches(2.85), Inches(6.3), Inches(0.38),
         size=13, bold=True, color=DARK)

statuts = [
    ("PENDING",    GRAY),
    ("ACCEPTED",   BLUE),
    ("PICKED_UP",  GREEN),
    ("IN_TRANSIT", RGBColor(234,88,12)),
    ("DELIVERED",  RGBColor(22,163,74)),
]
sx = Inches(6.6)
for sname, scol in statuts:
    add_rect(sl, sx, Inches(3.3), Inches(1.1), Inches(0.44), scol)
    add_text(sl, sname, sx, Inches(3.32), Inches(1.1), Inches(0.4),
             size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if sname != "DELIVERED":
        add_text(sl, "→", sx + Inches(1.1), Inches(3.3),
                 Inches(0.18), Inches(0.44),
                 size=12, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
    sx += Inches(1.28)

# Snapshot Pattern
add_rect(sl, Inches(6.6), Inches(4.05), Inches(6.3), Inches(0.04), BLUE)
add_rect(sl, Inches(6.6), Inches(4.09), Inches(6.3), Inches(0.9), LIGHT)
add_text(sl, "Snapshot Pattern (prix)",
         Inches(6.75), Inches(4.13), Inches(6.0), Inches(0.3),
         size=11, bold=True, color=BLUE)
add_text(sl, "unit_price_snapshot copié à la commande — protège contre les\nvariations de prix ultérieures.",
         Inches(6.75), Inches(4.43), Inches(6.0), Inches(0.5),
         size=10, color=BLACK)

# RLS Rating
add_rect(sl, Inches(6.6), Inches(5.15), Inches(6.3), Inches(0.04), GREEN)
add_rect(sl, Inches(6.6), Inches(5.19), Inches(6.3), Inches(0.6), LIGHT)
add_text(sl, "Contrainte évaluation (RLS)",
         Inches(6.75), Inches(5.23), Inches(6.0), Inches(0.28),
         size=11, bold=True, color=GREEN)
add_text(sl, "Un restaurateur ne peut évaluer qu'un agriculteur dont au moins 1 article est DELIVERED.",
         Inches(6.75), Inches(5.51), Inches(6.0), Inches(0.24),
         size=10, color=BLACK)

slide_footer(sl)


# =============================================================
# SLIDE 7 — MODULE SECONDSERVE
# =============================================================
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Module SecondServe — Anti-gaspillage alimentaire",
             "Offres surplus · Commandes citoyens · SSO Magic Link")

# Concept
add_text(sl, "Concept",
         Inches(0.4), Inches(1.42), Inches(6.0), Inches(0.38),
         size=14, bold=True, color=DARK)
concept = [
    "Les restaurants publient leurs surplus du jour à prix réduit",
    "Les citoyens filtrent par ville et réservent en ligne",
    "Paiement en espèces à la collecte (COD)",
]
add_bullet_block(sl, concept, Inches(0.4), Inches(1.86),
                 Inches(5.8), Inches(1.3), size=13, dot_color=RGBColor(234,88,12))

# Architecture
add_text(sl, "Architecture technique",
         Inches(0.4), Inches(3.25), Inches(6.0), Inches(0.38),
         size=14, bold=True, color=DARK)
arch = [
    "Application Vite 6 / React 19 indépendante",
    "Partagée sur le même projet Supabase que VitaChain",
    "Tables préfixées ss_",
    "Offres visibles uniquement si KYC approuvé + commerce_type défini",
    "Commande atomique via RPC ss_place_order() — contrôle de stock intégré",
]
add_bullet_block(sl, arch, Inches(0.4), Inches(3.68),
                 Inches(5.8), Inches(2.3), size=12)

# SSO diagram
add_text(sl, "SSO inter-application (Magic Link)",
         Inches(6.6), Inches(1.42), Inches(6.3), Inches(0.38),
         size=14, bold=True, color=DARK)

add_rect(sl, Inches(6.6), Inches(1.9), Inches(2.8), Inches(0.65), BLUE)
add_text(sl, "VitaChain",
         Inches(6.6), Inches(1.92), Inches(2.8), Inches(0.61),
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(sl, Inches(9.65), Inches(1.9), Inches(3.25), Inches(0.65),
         RGBColor(234, 88, 12))
add_text(sl, "SecondServe",
         Inches(9.65), Inches(1.92), Inches(3.25), Inches(0.61),
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(sl, Inches(9.45), Inches(2.07), Inches(0.2), Inches(0.3), GRAY)
add_text(sl, "OTP Supabase à usage unique",
         Inches(6.6), Inches(2.65), Inches(6.3), Inches(0.35),
         size=10, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

add_text(sl, "Même projet Supabase — tables ss_",
         Inches(6.6), Inches(2.98), Inches(6.3), Inches(0.3),
         size=10, color=GRAY, align=PP_ALIGN.CENTER)

# Étapes SSO
add_text(sl, "Flux SSO détaillé",
         Inches(6.6), Inches(3.42), Inches(6.3), Inches(0.35),
         size=13, bold=True, color=DARK)
sso_steps = [
    "1.  L'utilisateur clique « Accéder à SecondServe » dans VitaChain",
    "2.  Le backend génère un OTP Supabase à usage unique",
    "3.  Redirection vers SecondServe avec le token dans l'URL",
    "4.  SecondServe échange le token contre une session valide",
    "5.  Utilisateur connecté — sans ressaisie d'identifiants",
]
add_bullet_block(sl, sso_steps, Inches(6.6), Inches(3.82),
                 Inches(6.3), Inches(2.35), size=11,
                 dot_color=RGBColor(234, 88, 12))

slide_footer(sl)


# =============================================================
# SLIDE 8 — STACK TECHNIQUE
# =============================================================
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Stack Technique",
             "Technologies · Langages · Services Cloud · DevOps")

sections = [
    ("Backend & Base de données", BLUE, Inches(0.4), [
        "FastAPI 0.115 (Python 3.12) + Pydantic v2",
        "PostgreSQL 17 : 50 migrations · 12 tables · 5 vues · 8 ENUM",
        "Workers asynchrones : alertes, IA, offline, expiration",
        "ESP32 firmware C++ (Arduino IDE)",
    ]),
    ("Frontend", GREEN, Inches(0.4), [
        "Next.js 15.1 (RSC + SSR) + TypeScript 5 + Tailwind CSS",
        "Vite 6 + React 19 + React Router v7 (SecondServe)",
        "Leaflet (cartes parcelles GeoJSON)",
        "Recharts (graphiques télémétrie temps réel)",
    ]),
    ("Services Cloud & APIs", RGBColor(109,40,217), Inches(6.8), [
        "Supabase (Auth JWT · Storage S3 · Realtime)",
        "Google Gemini API (diagnostics agronomiques IA)",
        "OpenWeatherMap (prévisions météo parcelles)",
        "Copernicus / Sentinel-2 (NDVI satellite)",
        "Brevo (emails transactionnels)",
        "Backblaze B2 (sauvegardes nightly via rclone)",
    ]),
    ("DevOps & Qualité", RGBColor(234,88,12), Inches(6.8), [
        "Docker Compose (6 services conteneurisés)",
        "NGINX 1.27 + Certbot / Let's Encrypt",
        "GitHub Actions CI/CD (5 jobs parallèles)",
        "Ruff + ESLint + pgTAP + Locust (tests de charge)",
    ]),
]

row_y = [Inches(1.42), Inches(3.85)]
col_x = [Inches(0.4), Inches(6.8)]
col_w = Inches(6.1)

for idx, (title, col, cx, items) in enumerate(sections):
    ry = row_y[idx // 2]
    cx = col_x[idx % 2]
    add_rect(sl, cx, ry, col_w, Inches(0.38), col)
    add_text(sl, "  " + title, cx, ry + Inches(0.02),
             col_w, Inches(0.34), size=12, bold=True, color=WHITE)
    add_bullet_block(sl, items, cx, ry + Inches(0.42),
                     col_w, Inches(2.1), size=11, dot_color=col)

slide_footer(sl)


# =============================================================
# SLIDE 9 — MÉTHODOLOGIE SCRUM & CHIFFRES
# =============================================================
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Méthodologie Scrum & Résultats",
             "11 sprints de 2 semaines · Story-driven · CI/CD · ~460 heures")

# Tableau sprints
add_text(sl, "Tableau des sprints réalisés",
         Inches(0.4), Inches(1.42), Inches(6.0), Inches(0.38),
         size=13, bold=True, color=DARK)

sprint_rows = [
    ("Sprint 0",     "Infrastructure & DevOps"),
    ("Sprint 1",     "Authentification & KYC"),
    ("Sprints 2–4",  "Katara IoT  (6 semaines)"),
    ("Sprints 5–7",  "FarMarket   (6 semaines)"),
    ("Sprint 8",     "SecondServe"),
    ("Sprint 9",     "Administration"),
    ("Sprint 10",    "Tests & Livraison"),
]
add_rect(sl, Inches(0.4), Inches(1.88), Inches(6.0), Inches(0.38), BLUE)
add_text(sl, "Sprint", Inches(0.45), Inches(1.9), Inches(1.5), Inches(0.34),
         size=10, bold=True, color=WHITE)
add_text(sl, "Livrable", Inches(2.0), Inches(1.9), Inches(4.3), Inches(0.34),
         size=10, bold=True, color=WHITE)

for ri, (s, l) in enumerate(sprint_rows):
    bg = LIGHT if ri % 2 == 0 else WHITE
    y = Inches(2.26) + ri * Inches(0.46)
    add_rect(sl, Inches(0.4), y, Inches(6.0), Inches(0.46), bg)
    add_text(sl, s, Inches(0.5), y + Inches(0.07),
             Inches(1.4), Inches(0.32), size=11, bold=True, color=BLUE)
    add_text(sl, l, Inches(2.0), y + Inches(0.07),
             Inches(4.3), Inches(0.32), size=11, color=BLACK)

# KPIs
add_text(sl, "Chiffres clés du projet",
         Inches(6.8), Inches(1.42), Inches(6.0), Inches(0.38),
         size=13, bold=True, color=DARK)

kpis = [
    ("≈ 460 h",  "de développement",      BLUE,                  Inches(6.8),  Inches(1.9)),
    ("50+",      "migrations SQL",         GREEN,                 Inches(9.15), Inches(1.9)),
    ("11.5",     "mois de travail total",  RGBColor(109,40,217),  Inches(11.5), Inches(1.9)),
    ("≈ 60",     "endpoints REST API",     RGBColor(234,88,12),   Inches(6.8),  Inches(3.15)),
    ("4",        "modules livrés",         GREEN,                 Inches(9.15), Inches(3.15)),
    ("22",       "politiques RLS",         BLUE,                  Inches(11.5), Inches(3.15)),
]
for val, lbl, col, x, y in kpis:
    kpi_badge(sl, val, lbl, x, y, color=col)

# Scrum principles
add_text(sl, "Principes Scrum appliqués",
         Inches(6.8), Inches(4.4), Inches(6.0), Inches(0.35),
         size=12, bold=True, color=DARK)
principles = [
    "Développement itératif — livraison à chaque fin de sprint",
    "User stories avec critères d'acceptation dans docs/stories/",
    "Pre-commit hooks : Ruff, ESLint, shellcheck, AUTH-05",
    "Tests automatisés : pytest, pgTAP, tests E2E",
]
add_bullet_block(sl, principles, Inches(6.8), Inches(4.8),
                 Inches(6.0), Inches(1.8), size=11)

slide_footer(sl)


# =============================================================
# SLIDE 10 — CONCLUSION & PERSPECTIVES
# =============================================================
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Conclusion & Perspectives",
             "Plateforme production-ready · 4 modules · Déployée · Évolutive")

# Objectifs atteints
add_text(sl, "Objectifs atteints",
         Inches(0.4), Inches(1.42), Inches(6.0), Inches(0.38),
         size=14, bold=True, color=DARK)
achieved = [
    "Architecture full-stack production-ready",
    "KYC + JWT multi-rôle + RLS 22 politiques",
    "Katara IoT : ingestion < 50 ms (p50)",
    "FarMarket : cycle complet commandes B2B + anonymisation",
    "SecondServe : SSO + lutte anti-gaspillage",
    "Administration centralisée (KYC, users, modération)",
    "CI/CD + monitoring + sauvegardes automatisées",
]
add_bullet_block(sl, achieved, Inches(0.4), Inches(1.88),
                 Inches(5.8), Inches(3.0), size=12, dot_color=GREEN)

# Badge production
add_rect(sl, Inches(0.4), Inches(5.1), Inches(5.8), Inches(0.65), BLUE)
add_text(sl, "VitaChain n'est pas un prototype — elle est déployée avec\n"
             "HTTPS · CI/CD · sauvegardes nightly · monitoring actif",
         Inches(0.5), Inches(5.14), Inches(5.6), Inches(0.57),
         size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Perspectives
add_text(sl, "Perspectives",
         Inches(6.6), Inches(1.42), Inches(6.3), Inches(0.38),
         size=14, bold=True, color=DARK)

persp_sections = [
    ("Court terme", BLUE, [
        "Module BotaBa9a (surveillance chaîne du froid)",
        "Notifications push web (Web Push API)",
        "Application mobile native (React Native)",
    ]),
    ("Moyen terme", GREEN, [
        "Passerelle de paiement (CMI / PayZone Maroc)",
        "ML avancé : prédiction de rendement agricole",
        "Traçabilité sur la chaîne FarMarket",
    ]),
    ("Long terme", RGBColor(109,40,217), [
        "Marketplace de données agronomiques anonymisées",
        "Intégration API MAPMDREF (Ministère Agriculture)",
        "Réseau de capteurs IoT communautaires mutualisés",
    ]),
]
yp = Inches(1.88)
for ptitle, pcol, pitems in persp_sections:
    add_rect(sl, Inches(6.6), yp, Inches(6.3), Inches(0.3), pcol)
    add_text(sl, "  " + ptitle,
             Inches(6.6), yp + Inches(0.02), Inches(6.3), Inches(0.26),
             size=11, bold=True, color=WHITE)
    add_bullet_block(sl, pitems, Inches(6.6), yp + Inches(0.32),
                     Inches(6.3), Inches(1.1),
                     size=11, dot_color=pcol)
    yp += Inches(1.58)

slide_footer(sl)


# =============================================================
# SLIDE 11 — MERCI / QUESTIONS
# =============================================================
sl = prs.slides.add_slide(BLANK)

# Fond blanc, bandes
add_rect(sl, 0, 0, W, Inches(1.4), BLUE)
add_rect(sl, 0, Inches(1.4), W, Inches(0.07), GREEN)
add_rect(sl, 0, H - Inches(1.15), W, Inches(0.07), GREEN)
add_rect(sl, 0, H - Inches(1.08), W, Inches(1.08), BLUE)

add_text(sl, "VitaChain — PFE 2026",
         Inches(0.5), Inches(0.3), W - Inches(1.0), Inches(0.7),
         size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(sl, "Merci de votre attention",
         Inches(0.5), Inches(1.9), W - Inches(1.0), Inches(1.1),
         size=42, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

add_text(sl, "Nous sommes disponibles pour vos questions",
         Inches(0.5), Inches(3.1), W - Inches(1.0), Inches(0.6),
         size=18, color=GRAY, align=PP_ALIGN.CENTER)

# Filet décoratif
add_rect(sl, Inches(3.0),  Inches(3.85), Inches(2.8), Inches(0.06), BLUE)
add_rect(sl, Inches(5.8),  Inches(3.85), Inches(1.7), Inches(0.06), GREEN)
add_rect(sl, Inches(7.5),  Inches(3.85), Inches(2.8), Inches(0.06), BLUE)

# Noms
add_text(sl, "Badre Saad          Kodoussi Mohammed          El Karmi Yasser",
         Inches(0.5), Inches(4.25), W - Inches(1.0), Inches(0.5),
         size=15, bold=True, color=DARK, align=PP_ALIGN.CENTER)

add_text(sl,
         "Encadré par : Mme Ouchra Hafssa  ·  Mme Achtaich Khadija  ·  M. Ait Daoud Mohammed",
         Inches(0.5), Inches(4.85), W - Inches(1.0), Inches(0.4),
         size=12, color=GRAY, align=PP_ALIGN.CENTER)

add_text(sl,
         "Faculté des Sciences Ben M'Sik  —  Département Informatique  —  DIARS  —  2025 / 2026",
         Inches(0.5), Inches(5.3), W - Inches(1.0), Inches(0.4),
         size=11, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

add_text(sl, "Soutenu le 19 juin 2026",
         Inches(0.5), Inches(5.75), W - Inches(1.0), Inches(0.35),
         size=11, color=GRAY, align=PP_ALIGN.CENTER)


# =============================================================
# SAUVEGARDE
# =============================================================
OUT = "docs/presentation_vitachain.pptx"
prs.save(OUT)
print(f"OK - Presentation generee : {OUT}")
print(f"   Slides : {len(prs.slides)}")
